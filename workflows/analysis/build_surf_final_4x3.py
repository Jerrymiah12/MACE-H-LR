#!/usr/bin/env python3
"""Build Jeremiah Bailey's 15-minute, general-audience SURF final deck."""
from __future__ import annotations

import json
import sys
import tempfile
from pathlib import Path
from zipfile import ZipFile

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

ROOT = Path(__file__).resolve().parents[2]
HERE = ROOT / "results" / "learned_response"

from workflows.analysis.build_project_update_presentation import (
    BLUE, GRAY, GREEN, INK, LIGHT, LINE, MID, MINT, NAVY, ORANGE,
    PALE_BLUE, PALE_ORANGE, RED, TEAL, WHITE, add_notes,
    add_picture_fit, box, line, metric, new_slide, rich_text, status_chip,
    text,
)
import workflows.analysis.build_project_update_presentation as theme


OUTPUT = HERE / "Jeremiah_Bailey_SURF_Final_Presentation_4x3_2026.pptx"
NOTES_OUTPUT = HERE / "Jeremiah_Bailey_SURF_Final_Presentation_4x3_2026_speaker_notes.md"
QR_DOCX = ROOT / "docs" / "archive" / "Audience Feedback QR Code.docx"
W, H = 10.0, 7.5


def title(slide, kicker, headline, number):
    text(slide, kicker.upper(), 0.58, 0.27, 8.8, 0.20, size=10,
         color=TEAL, bold=True)
    size = 24 if len(headline) > 52 else 29
    text(slide, headline, 0.58, 0.55, 8.82, 0.67, size=size,
         color=NAVY, font="Cambria", bold=True)
    line(slide, 0.58, 1.35, 9.42, 1.35, color=LINE, width=0.9)
    text(slide, str(number), 9.10, 0.30, 0.30, 0.17, size=8.5,
         color=GRAY, align=PP_ALIGN.RIGHT)


def footer(slide, number):
    text(slide, "JEREMIAH BAILEY  ·  SURF 2026  ·  BERNARDI GROUP",
         0.58, 7.20, 5.75, 0.13, size=7.4, color=GRAY)
    text(slide, str(number), 9.08, 7.20, 0.32, 0.13, size=7.4,
         color=GRAY, align=PP_ALIGN.RIGHT)


def picture(slide, name, x, y, w, h):
    return add_picture_fit(slide, HERE / name, x, y, w, h)


def atom(slide, x, y, size, fill, label=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y),
                                   Inches(size), Inches(size))
    shape.fill.solid()
    shape.fill.fore_color.rgb = theme.rgb(fill)
    shape.line.color.rgb = theme.rgb(WHITE)
    shape.line.width = theme.Pt(1.2)
    if label:
        text(slide, label, x, y + size * 0.28, size, size * 0.30,
             size=9, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    return shape


def comparison_bar(slide, label, learned, baseline, x, y, w, color, unit):
    text(slide, label, x, y, w, 0.24, size=12, color=NAVY, bold=True)
    maximum = max(learned, baseline)
    rows = (("Constant baseline", baseline, GRAY),
            ("Learned model", learned, color))
    for index, (name, value, shade) in enumerate(rows):
        yy = y + 0.48 + index * 0.68
        text(slide, name, x, yy, 1.35, 0.18, size=8.5, color=shade,
             bold=(index == 1))
        box(slide, x + 1.43, yy, w - 1.43, 0.20, fill=LINE, radius=True)
        box(slide, x + 1.43, yy, (w - 1.43) * value / maximum, 0.20,
            fill=shade, radius=True)
        text(slide, f"{value:.6f}{unit}", x + 1.43, yy + 0.27,
             w - 1.43, 0.17, size=8.2, color=shade,
             align=PP_ALIGN.RIGHT)


def notes(slide, heading, body):
    add_notes(slide, heading, body)


def build():
    metrics_data = json.loads((HERE / "metrics.json").read_text())
    h = metrics_data["hamiltonian"]
    epc = metrics_data["epc"]
    tensors = metrics_data["tensor_targets"]
    born = tensors["born"]
    epsilon = tensors["epsilon"]
    baseline = tensors["constant_baseline"]

    theme.SLIDE_NOTES.clear()
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    with tempfile.TemporaryDirectory(prefix="surf_final_assets_") as tmp:
        tmp_dir = Path(tmp)
        with ZipFile(QR_DOCX) as archive:
            qr = tmp_dir / "feedback_qr.png"
            logo = tmp_dir / "caltech_logo.png"
            qr.write_bytes(archive.read("word/media/image1.png"))
            logo.write_bytes(archive.read("word/media/image2.png"))

        # 1 — title
        s = new_slide(prs, NAVY)
        box(s, 0, 0, 0.20, H, fill=MINT)
        status_chip(s, "SURF final presentation", 0.74, 0.68, 1.92,
                    color=MINT, text_color=NAVY)
        text(s, "Teaching machines how atoms\nand electrons move together",
             0.74, 1.42, 8.60, 1.38, size=33, color=WHITE,
             font="Cambria", bold=True)
        text(s, "A physics-guided machine-learning model for polar materials",
             0.76, 3.05, 8.30, 0.52, size=19, color=PALE_BLUE,
             font="Cambria", bold=True)
        line(s, 0.76, 3.91, 8.92, 3.91, color=TEAL, width=1.0)
        text(s, "Jeremiah Bailey", 0.76, 4.34, 8.20, 0.42, size=21,
             color=WHITE, font="Cambria", bold=True)
        text(s, "Bernardi Group · California Institute of Technology",
             0.76, 4.90, 8.20, 0.30, size=12.5, color=MID)
        text(s, "SURF 2026", 0.76, 5.33, 8.20, 0.25, size=11,
             color=MINT, bold=True)
        add_picture_fit(s, logo, 7.66, 5.65, 1.55, 0.84, border=None)
        notes(s, "Teaching machines how atoms and electrons move together", """
0:00–0:35. Hello, my name is Jeremiah Bailey. This summer I worked in the
Bernardi Group on a question at the intersection of materials physics and
machine learning: can we teach a model not only the electronic state of a
material, but also how that state changes when the atoms move? I will focus on
magnesium oxide, a simple polar crystal that makes this challenge visible.
""")

        # 2 — why it matters
        s = new_slide(prs)
        title(s, "Why this matters", "Atoms vibrate—and electrons must respond", 2)
        box(s, 0.62, 1.58, 4.16, 4.86, fill=NAVY, radius=True)
        text(s, "A crystal is not static", 0.94, 1.91, 3.50, 0.36,
             size=21, color=WHITE, font="Cambria", bold=True,
             align=PP_ALIGN.CENTER)
        # Simple MgO lattice illustration.
        x0, y0 = 1.23, 2.70
        for row in range(3):
            for col in range(3):
                xx, yy = x0 + col * 1.02, y0 + row * 0.87
                fill = ORANGE if (row + col) % 2 == 0 else BLUE
                atom(s, xx, yy, 0.44, fill, "Mg" if fill == ORANGE else "O")
                if col < 2:
                    line(s, xx + 0.44, yy + 0.22, xx + 1.02, yy + 0.22,
                         color=MID, width=1.2)
                if row < 2:
                    line(s, xx + 0.22, yy + 0.44, xx + 0.22, yy + 0.87,
                         color=MID, width=1.2)
        text(s, "atomic vibration", 1.20, 5.52, 1.50, 0.23, size=9.5,
             color=MINT, bold=True)
        line(s, 2.40, 5.65, 3.45, 5.65, color=MINT, width=2.3, arrow=True)
        text(s, "electronic response", 3.05, 5.52, 1.40, 0.23, size=9.5,
             color=MINT, bold=True, align=PP_ALIGN.RIGHT)

        text(s, "This interaction helps determine:", 5.24, 1.74, 4.06, 0.34,
             size=18, color=NAVY, font="Cambria", bold=True)
        items = [
            ("How electrical energy is lost", "resistance and heat", ORANGE),
            ("How carriers exchange energy", "relaxation and transport", BLUE),
            ("How materials respond to light", "optical and polar behavior", GREEN),
        ]
        for i, (head, body, color) in enumerate(items):
            y = 2.42 + i * 1.18
            box(s, 5.24, y, 0.42, 0.42, fill=color, radius=True)
            text(s, str(i + 1), 5.24, y + 0.06, 0.42, 0.18, size=10,
                 color=WHITE, bold=True, align=PP_ALIGN.CENTER)
            text(s, head, 5.90, y - 0.01, 3.26, 0.30, size=14,
                 color=NAVY, bold=True)
            text(s, body, 5.90, y + 0.39, 3.26, 0.23, size=10.5,
                 color=GRAY)
        box(s, 5.24, 6.08, 4.02, 0.48, fill=LIGHT, line=LINE, radius=True)
        text(s, "Goal: predict the response without repeating an expensive quantum calculation for every motion.",
             5.47, 6.19, 3.56, 0.24, size=9.7, color=NAVY, bold=True,
             align=PP_ALIGN.CENTER)
        footer(s, 2)
        notes(s, "Atoms vibrate—and electrons must respond", """
0:35–1:35. Atoms in a solid are always moving. When they move, the electrons
rearrange. That interaction affects resistance, energy relaxation, heat, and
optical behavior. Quantum-mechanical calculations can describe it accurately,
but repeating those calculations for every possible atomic motion is costly.
The opportunity is to learn this response from examples while keeping the
relevant physics in the model.
""")

        # 3 — roadmap
        s = new_slide(prs)
        title(s, "Roadmap", "Four questions guide the talk", 3)
        questions = [
            ("1", "What are we predicting?", "A compact electronic description and how it changes when atoms move.", BLUE),
            ("2", "Why is magnesium oxide hard?", "Its positive and negative ions create a long-range electric response.", ORANGE),
            ("3", "What did I build and test?", "A shared model with learned charge and screening outputs.", GREEN),
            ("4", "What do the results say?", "The new outputs learn real signal, but have not changed the final response yet.", TEAL),
        ]
        for i, (num, head, body, color) in enumerate(questions):
            y = 1.67 + i * 1.27
            box(s, 0.82, y, 0.52, 0.52, fill=color, radius=True)
            text(s, num, 0.82, y + 0.08, 0.52, 0.22, size=13,
                 color=WHITE, bold=True, align=PP_ALIGN.CENTER)
            text(s, head, 1.62, y - 0.01, 3.10, 0.30, size=17,
                 color=NAVY, font="Cambria", bold=True)
            text(s, body, 4.80, y - 0.03, 4.25, 0.62, size=11.5,
                 color=GRAY, valign=MSO_ANCHOR.MIDDLE)
            if i < 3:
                line(s, 1.08, y + 0.54, 1.08, y + 1.17, color=LINE,
                     width=2.0)
        footer(s, 3)
        notes(s, "Four questions guide the talk", """
1:35–2:05. I will first define the quantities in plain language. Then I will
explain why polar magnesium oxide needs special treatment, show the model and
training protocol I built, and finish with what the data support so far and
what experiment should come next.
""")

        # 4 — accessible background
        s = new_slide(prs)
        title(s, "Background", "Three ideas connect structure to response", 4)
        cards = [
            (0.62, "1", "DFT", "Physics-based reference", "Solves an approximate quantum-mechanical problem for one atomic structure.", NAVY),
            (3.57, "2", "Hamiltonian", "Electronic rulebook", "A matrix that summarizes the energies and interactions available to electrons.", BLUE),
            (6.52, "3", "Electron–phonon coupling", "The local slope", "Measures how the Hamiltonian changes when atoms are displaced.", ORANGE),
        ]
        for x, num, head, sub, body, color in cards:
            box(s, x, 1.70, 2.66, 4.55, fill=LIGHT, line=LINE, radius=True)
            box(s, x + 0.20, 1.95, 0.44, 0.44, fill=color, radius=True)
            text(s, num, x + 0.20, 2.02, 0.44, 0.18, size=11,
                 color=WHITE, bold=True, align=PP_ALIGN.CENTER)
            text(s, head, x + 0.22, 2.62, 2.22, 0.42, size=21,
                 color=NAVY, font="Cambria", bold=True,
                 align=PP_ALIGN.CENTER)
            text(s, sub.upper(), x + 0.22, 3.22, 2.22, 0.20, size=8.5,
                 color=color, bold=True, align=PP_ALIGN.CENTER)
            text(s, body, x + 0.30, 3.74, 2.06, 1.16, size=11.2,
                 color=GRAY, align=PP_ALIGN.CENTER,
                 valign=MSO_ANCHOR.MIDDLE)
            if num == "3":
                text(s, "change in electrons\n──────────────\nchange in atoms",
                     x + 0.52, 5.14, 1.62, 0.72, size=10, color=ORANGE,
                     bold=True, align=PP_ALIGN.CENTER)
            elif num == "2":
                # Small matrix icon.
                for rr in range(3):
                    for cc in range(3):
                        box(s, x + 0.87 + cc * 0.28, 5.15 + rr * 0.25,
                            0.18, 0.15, fill=BLUE if rr == cc else PALE_BLUE,
                            radius=True)
            else:
                text(s, "accurate · expensive", x + 0.45, 5.34, 1.76, 0.22,
                     size=10, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        box(s, 1.04, 6.52, 7.92, 0.38, fill=NAVY, radius=True)
        text(s, "Machine learning approximates the electronic rulebook; the response tests whether it learned the right behavior between examples.",
             1.28, 6.61, 7.44, 0.20, size=9.8, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER)
        footer(s, 4)
        notes(s, "Three ideas connect structure to response", """
2:05–3:20. DFT is my physics-based reference: accurate, but expensive for
large numbers of structures. The Hamiltonian is a compact matrix that acts
like an electronic rulebook. Machine learning can predict that matrix from
atomic positions. Electron–phonon coupling is more demanding: it is the slope
of that prediction when an atom moves. A model can match many values and still
learn the wrong slope between them.
""")

        # 5 — MgO challenge and four-way comparison
        s = new_slide(prs)
        title(s, "The challenge", "Polar crystals communicate over long distances", 5)
        box(s, 0.60, 1.60, 3.15, 4.95, fill=NAVY, radius=True)
        text(s, "Magnesium oxide", 0.86, 1.93, 2.62, 0.37, size=20,
             color=WHITE, font="Cambria", bold=True,
             align=PP_ALIGN.CENTER)
        atom(s, 1.12, 2.77, 0.76, ORANGE, "Mg²⁺")
        atom(s, 2.47, 2.77, 0.76, BLUE, "O²⁻")
        line(s, 1.92, 3.15, 2.41, 3.15, color=MINT, width=2.8, arrow=True)
        text(s, "A displacement creates an electric field that reaches beyond nearby atoms.",
             0.92, 4.05, 2.50, 0.91, size=12, color=PALE_BLUE,
             bold=True, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE)
        text(s, "Local chemistry + long-range electrostatics",
             0.92, 5.46, 2.50, 0.42, size=10.5, color=MINT, bold=True,
             align=PP_ALIGN.CENTER)

        text(s, "Four outputs in the comparison", 4.16, 1.73, 5.08, 0.32,
             size=18, color=NAVY, font="Cambria", bold=True)
        rows = [
            ("Actual DFT", "physics reference", NAVY),
            ("Direct Full-H", "learn the complete electronic matrix", ORANGE),
            ("LR-corrected SR", "learn local part + fixed long-range physics", BLUE),
            ("SR + predicted Z* + ε∞", "learn local part + learned charge and screening", GREEN),
        ]
        for i, (name, body, color) in enumerate(rows):
            y = 2.28 + i * 0.91
            box(s, 4.18, y, 0.42, 0.42, fill=color, radius=True)
            text(s, str(i + 1), 4.18, y + 0.06, 0.42, 0.18, size=10,
                 color=WHITE, bold=True, align=PP_ALIGN.CENTER)
            text(s, name, 4.85, y - 0.01, 2.18, 0.27, size=13.5,
                 color=NAVY, bold=True)
            text(s, body, 7.10, y - 0.01, 2.15, 0.43, size=9.5,
                 color=GRAY)
        box(s, 4.18, 6.13, 5.08, 0.43, fill=LIGHT, line=LINE, radius=True)
        text(s, "Z* = effective charge · ε∞ = electronic screening",
             4.42, 6.25, 4.60, 0.18, size=9.8, color=NAVY, bold=True,
             align=PP_ALIGN.CENTER)
        footer(s, 5)
        notes(s, "Polar crystals communicate over long distances", """
3:20–4:25. Magnesium and oxygen carry opposite ionic character. Moving one
atom creates an electric response that is not purely local. I compare four
things: the DFT reference, a model trained directly on the complete
Hamiltonian, a short-range model corrected with fixed reference physics, and
my new model whose effective charges and electronic screening are predicted
from the current geometry.
""")

        # 6 — what I built
        s = new_slide(prs)
        title(s, "What I built", "One structural encoder feeds three physical predictions", 6)
        box(s, 0.62, 2.29, 1.28, 1.10, fill=LIGHT, line=LINE, radius=True)
        text(s, "Atomic\nstructure", 0.77, 2.53, 0.98, 0.48, size=15,
             color=NAVY, font="Cambria", bold=True,
             align=PP_ALIGN.CENTER)
        line(s, 1.94, 2.84, 2.30, 2.84, color=TEAL, width=2.5, arrow=True)
        box(s, 2.37, 1.98, 1.53, 1.72, fill=NAVY, radius=True)
        text(s, "Shared\nencoder", 2.55, 2.37, 1.17, 0.56, size=17,
             color=WHITE, font="Cambria", bold=True,
             align=PP_ALIGN.CENTER)
        text(s, "learns reusable\nstructure features", 2.55, 3.09, 1.17, 0.35,
             size=8.7, color=MID, align=PP_ALIGN.CENTER)
        heads = [
            (4.48, 1.62, "Short-range H", "local electronic rulebook", BLUE, "E7F2F8"),
            (4.48, 2.75, "Born charge Z*", "how motion polarizes the crystal", ORANGE, PALE_ORANGE),
            (4.48, 3.88, "Screening ε∞", "how electrons weaken the field", GREEN, "E4F4EE"),
        ]
        for x, y, head, body, color, fill in heads:
            line(s, 3.94, 2.84, x - 0.08, y + 0.45, color=TEAL,
                 width=1.7, arrow=True)
            box(s, x, y, 2.10, 0.90, fill=fill, line=LINE, radius=True)
            text(s, head, x + 0.12, y + 0.14, 1.86, 0.27, size=14,
                 color=color, font="Cambria", bold=True,
                 align=PP_ALIGN.CENTER)
            text(s, body, x + 0.13, y + 0.52, 1.84, 0.20, size=8.8,
                 color=GRAY, align=PP_ALIGN.CENTER)
        line(s, 6.67, 2.84, 7.05, 2.84, color=TEAL, width=2.4, arrow=True)
        box(s, 7.12, 2.10, 2.18, 1.48, fill=NAVY, radius=True)
        text(s, "Reconstruct\ntotal H", 7.34, 2.39, 1.74, 0.52, size=18,
             color=WHITE, font="Cambria", bold=True,
             align=PP_ALIGN.CENTER)
        text(s, "local ML + analytic\nlong-range physics", 7.34, 3.06, 1.74, 0.35,
             size=8.8, color=MINT, align=PP_ALIGN.CENTER)
        box(s, 0.82, 5.43, 8.36, 0.77, fill=LIGHT, line=LINE, radius=True)
        rich_text(s, [
            ("Built-in safeguards:  ", NAVY, True, 11),
            ("charge balance", ORANGE, True, 11),
            ("  ·  ", GRAY, False, 11),
            ("positive screening", GREEN, True, 11),
            ("  ·  ", GRAY, False, 11),
            ("no direct full-H shortcut", BLUE, True, 11),
        ], 1.12, 5.68, 7.76, 0.27, align=PP_ALIGN.CENTER)
        footer(s, 6)
        notes(s, "One structural encoder feeds three physical predictions", """
4:25–5:25. I built one shared equivariant encoder, meaning its outputs transform
correctly when the crystal is rotated. It feeds three predictions: the local
Hamiltonian, an effective charge for every atom, and one screening tensor for
the structure. Those pieces are combined with an analytic long-range formula.
The model enforces overall charge balance and positive screening by design.
""")

        # 7 — data protocol
        s = new_slide(prs)
        title(s, "How I tested it", "Separate data answer separate questions", 7)
        sections = [
            (0.64, "Hamiltonian data", "367 structures", "Teach and monitor the electronic matrix", BLUE),
            (3.68, "Tensor training", "10 structures", "Teach effective charge and screening", ORANGE),
            (6.72, "Locked tensor test", "5 structures", "Open once after every model choice is frozen", GREEN),
        ]
        for x, head, count, body, color in sections:
            box(s, x, 1.72, 2.64, 3.43, fill=LIGHT, line=LINE, radius=True)
            status_chip(s, count, x + 0.74, 1.98, 1.17, color=color)
            text(s, head, x + 0.20, 2.60, 2.24, 0.38, size=18,
                 color=NAVY, font="Cambria", bold=True,
                 align=PP_ALIGN.CENTER)
            text(s, body, x + 0.30, 3.38, 2.04, 0.83, size=11,
                 color=GRAY, align=PP_ALIGN.CENTER,
                 valign=MSO_ANCHOR.MIDDLE)
            if x < 6:
                line(s, x + 2.66, 3.42, x + 2.96, 3.42, color=TEAL,
                     width=2.2, arrow=True)
        box(s, 0.90, 5.56, 8.20, 0.74, fill=NAVY, radius=True)
        text(s, "Selection rule", 1.16, 5.76, 1.28, 0.22, size=11,
             color=MINT, bold=True)
        text(s, "Hamiltonian, charge, and screening gates had to pass for ten consecutive validation checks.",
             2.54, 5.73, 6.12, 0.28, size=10.5, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER)
        text(s, "The locked five were not used to select the architecture, learning rate, or checkpoint.",
             1.12, 6.52, 7.76, 0.27, size=9.5, color=GRAY,
             italic=True, align=PP_ALIGN.CENTER)
        footer(s, 7)
        notes(s, "Separate data answer separate questions", """
5:25–6:20. Hamiltonians are large, so the model retained all 367 short-range
training and validation structures. The response tensors are expensive DFT
labels: ten trained the heads, two guided wiring and stability, and five were
locked. I opened those five only after freezing the architecture, learning
rates, and checkpoint. The accepted model had to pass all three validation
gates ten times in a row.
""")

        # 8 — optimization
        s = new_slide(prs)
        title(s, "Training outcome", "The new tasks did not destabilize the original model", 8)
        picture(s, "01_training_and_multitask_gates.png", 0.48, 1.55, 9.04, 4.18)
        box(s, 0.78, 5.95, 8.44, 0.45, fill=NAVY, radius=True)
        text(s, "All three gates passed for ten consecutive epochs · selected checkpoint: epoch 9",
             1.00, 6.07, 8.00, 0.20, size=10.5, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER)
        text(s, "Read the center and right panels as stability checks; the detailed scales are less important than staying below each dashed threshold.",
             0.88, 6.60, 8.24, 0.26, size=8.8, color=GRAY,
             italic=True, align=PP_ALIGN.CENTER)
        footer(s, 8)
        notes(s, "The new tasks did not destabilize the original model", """
6:20–7:15. The left panel gives the longer training context. The center asks
whether adding new tasks damages the original short-range Hamiltonian; it
stays below the dashed limit. The right asks whether the charge and screening
heads beat simple baselines; both stay below one. After ten consecutive passes,
training stopped and I froze the checkpoint.
""")

        # 9 — locked tensor result
        s = new_slide(prs)
        title(s, "First result", "The learned tensors beat constant guesses on unseen structures", 9)
        comparison_bar(s, "Effective charge (Born tensor)", born["mae"],
                       baseline["born_mae"], 0.74, 1.73, 4.08, ORANGE, " e")
        comparison_bar(s, "Electronic screening (ε∞)", epsilon["mae"],
                       baseline["epsilon_mae"], 5.18, 1.73, 4.08, GREEN, "")
        metric(s, f"{100*(1-born['mae']/baseline['born_mae']):.0f}%",
               "LOWER CHARGE ERROR", 0.92, 4.18, 1.82,
               color=ORANGE, fill=PALE_ORANGE)
        metric(s, f"{100*(1-epsilon['mae']/baseline['epsilon_mae']):.0f}%",
               "LOWER SCREENING ERROR", 2.94, 4.18, 1.92,
               color=GREEN, fill="E4F4EE")
        metric(s, "balanced", "TOTAL EFFECTIVE CHARGE", 5.17, 4.18, 1.83,
               color=BLUE, fill="E7F2F8")
        metric(s, "positive", "SCREENING EIGENVALUES", 7.22, 4.18, 1.88,
               color=TEAL)
        box(s, 0.86, 5.44, 8.28, 0.66, fill=LIGHT, line=LINE, radius=True)
        text(s, "Interpretation: the heads learned geometry-dependent signal. Limitation: five test structures support a prototype result, not broad generalization.",
             1.12, 5.61, 7.76, 0.33, size=10.2, color=NAVY, bold=True,
             align=PP_ALIGN.CENTER)
        footer(s, 9)
        notes(s, "The learned tensors beat constant guesses on unseen structures", """
7:15–8:15. This is the first success criterion. On structures never used for
training or selection, the learned effective charges reduce mean absolute
error by about 26 percent compared with always predicting the training mean.
Screening error drops by about 73 percent. The outputs also satisfy the
physical constraints. This shows learned geometry-dependent signal, while the
small test set keeps the claim at prototype scale.
""")

        # 10 — H comparison
        s = new_slide(prs)
        title(s, "Second result", "For electronic values, direct Full-H still leads", 10)
        picture(s, "02_locked_hamiltonian_error.png", 0.48, 1.56, 7.14, 4.42)
        metric(s, f"{1e3*h['full']['mae']:.3f}", "DIRECT FULL-H MAE (meV)",
               7.80, 1.80, 1.60, color=ORANGE, fill=PALE_ORANGE)
        metric(s, f"{1e3*h['sr']['mae']:.3f}", "FIXED-PHYSICS SR (meV)",
               7.80, 3.03, 1.60, color=BLUE, fill="E7F2F8")
        metric(s, f"{1e3*h['tensor']['mae']:.3f}", "LEARNED-TENSOR SR (meV)",
               7.80, 4.26, 1.60, color=GREEN, fill="E4F4EE")
        text(s, "Lower is better", 7.80, 5.48, 1.60, 0.18, size=8.8,
             color=GRAY, bold=True, align=PP_ALIGN.CENTER)
        box(s, 0.86, 6.28, 8.28, 0.50, fill=NAVY, radius=True)
        text(s, "The learned tensors slightly improve the SR reconstruction—but do not close the value gap.",
             1.10, 6.41, 7.80, 0.22, size=10.2, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER)
        footer(s, 10)
        notes(s, "For electronic values, direct Full-H still leads", """
8:15–9:15. Next I compare matrix values on those same five structures. Zero is
the DFT reference. Direct full-H has the smallest mean error at 0.417
millielectron-volts. The fixed-physics SR result is 0.487, and learned tensors
slightly improve it to 0.485. So the composed model is complete and stable,
but the tensor heads do not close the value-accuracy gap.
""")

        # 11 — EPC overall
        s = new_slide(prs)
        title(s, "Third result", "For atomic-motion response, both SR models win", 11)
        picture(s, "04_epc_overall_error.png", 0.48, 1.58, 7.18, 4.37)
        metric(s, f"{100*epc['full']['relative_l2']:.0f}%", "DIRECT FULL-H ERROR",
               7.82, 1.87, 1.58, color=ORANGE, fill=PALE_ORANGE)
        metric(s, f"{100*epc['sr']['relative_l2']:.1f}%", "FIXED-PHYSICS SR ERROR",
               7.82, 3.15, 1.58, color=BLUE, fill="E7F2F8")
        metric(s, f"{100*epc['tensor']['relative_l2']:.1f}%", "LEARNED-TENSOR SR ERROR",
               7.82, 4.43, 1.58, color=GREEN, fill="E4F4EE")
        box(s, 0.86, 6.27, 8.28, 0.50, fill=NAVY, radius=True)
        text(s, "Matching values is not enough: the derivative exposes how the model behaves between examples.",
             1.10, 6.40, 7.80, 0.22, size=10.2, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER)
        footer(s, 11)
        notes(s, "For atomic-motion response, both SR models win", """
9:15–10:20. The ranking reverses for the response. Here lower is again better.
Direct full-H has about 92 percent relative error, while both short-range
pipelines are near 24 percent. This is the central lesson: a model can predict
Hamiltonian values well and still learn the wrong slope. Physics-guided target
design improves the derivative observable even when it does not win on values.
""")

        # 12 — q resolved
        s = new_slide(prs)
        title(s, "Where the difference appears", "Longer-wavelength labels separate value and response quality", 12)
        picture(s, "05_epc_q_resolved.png", 0.54, 1.56, 6.22, 4.87)
        box(s, 7.04, 1.79, 2.28, 4.31, fill=LIGHT, line=LINE, radius=True)
        text(s, "How to read this", 7.30, 2.05, 1.76, 0.30, size=17,
             color=NAVY, font="Cambria", bold=True,
             align=PP_ALIGN.CENTER)
        explanations = [
            ("q", "labels a vibration wavelength and direction", TEAL),
            ("Top", "overall coupling strength", BLUE),
            ("Bottom", "error relative to DFT", ORANGE),
        ]
        for i, (head, body, color) in enumerate(explanations):
            y = 2.62 + i * 0.86
            status_chip(s, head, 7.28, y, 0.70, color=color)
            text(s, body, 8.12, y + 0.01, 0.96, 0.42, size=8.8,
                 color=GRAY)
        box(s, 7.27, 5.29, 1.82, 0.53, fill=NAVY, radius=True)
        text(s, "Blue and green overlap; orange grows too strong away from q = 0.",
             7.43, 5.38, 1.50, 0.34, size=8.2, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER)
        footer(s, 12)
        notes(s, "Longer-wavelength labels separate value and response quality", """
10:20–11:25. q is simply a label for a vibration's wavelength and direction.
The top panel shows coupling strength; the bottom shows error against DFT. At
q equals zero, all models are much closer. Away from zero, direct full-H
becomes too strong. The blue fixed-tensor and green learned-tensor curves
overlap across this grid, which leads to the final interpretation.
""")

        # 13 — interpretation
        s = new_slide(prs)
        title(s, "What we learned", "The new heads work—but their EPC effect is below present sensitivity", 13)
        box(s, 0.62, 1.67, 4.20, 4.63, fill="E4F4EE", line=LINE, radius=True)
        status_chip(s, "Supported", 0.88, 1.94, 1.04, color=GREEN)
        supported = [
            "Charge and screening heads learn unseen variation.",
            "The reconstructed Hamiltonian is complete and physically constrained.",
            "Short-range target design improves EPC compared with direct Full-H.",
        ]
        for i, value in enumerate(supported):
            y = 2.61 + i * 0.99
            box(s, 0.91, y, 0.31, 0.31, fill=GREEN, radius=True)
            text(s, "✓", 0.91, y + 0.04, 0.31, 0.16, size=10,
                 color=WHITE, bold=True, align=PP_ALIGN.CENTER)
            text(s, value, 1.42, y - 0.02, 3.02, 0.66, size=11.2,
                 color=NAVY, bold=True)
        box(s, 5.18, 1.67, 4.20, 4.63, fill=PALE_ORANGE, line=LINE,
            radius=True)
        status_chip(s, "Not yet supported", 5.44, 1.94, 1.47, color=ORANGE)
        open_items = [
            "Geometry-dependent tensors improve EPC at 5×10⁻⁶ Å.",
            "Ten tensor-training structures generalize broadly.",
            "A 2×2×2 q grid resolves the long-wavelength polar limit.",
        ]
        for i, value in enumerate(open_items):
            y = 2.61 + i * 0.99
            box(s, 5.47, y, 0.31, 0.31, fill=ORANGE, radius=True)
            text(s, "?", 5.47, y + 0.04, 0.31, 0.16, size=10,
                 color=WHITE, bold=True, align=PP_ALIGN.CENTER)
            text(s, value, 5.98, y - 0.02, 3.02, 0.66, size=11.2,
                 color=NAVY, bold=True)
        text(s, "Frozen and geometry-dependent tensor EPC differ by only 4.3×10⁻¹⁰ eV/Å on average.",
             1.16, 6.54, 7.68, 0.28, size=9.4, color=GRAY,
             italic=True, align=PP_ALIGN.CENTER)
        footer(s, 13)
        notes(s, "The new heads work—but their EPC effect is below present sensitivity", """
11:25–12:30. The heads work as predictors, and the composed reconstruction is
correct. But the learned geometry dependence changes EPC by only about four
times ten to the minus ten electron-volts per ångström on average at this
finite-difference step. That does not prove the tensors are physically
irrelevant. It means this experiment cannot yet separate their effect from
near-neutrality. The dataset and q grid also limit the scope of the claim.
""")

        # 14 — next steps
        s = new_slide(prs)
        title(s, "Next steps", "Make the tensor effect measurable before scaling up", 14)
        next_steps = [
            ("1", "Displacement sensitivity", "Repeat the response calculation over larger atomic steps to find where geometry-dependent tensors rise above numerical uncertainty.", TEAL),
            ("2", "More informative labels", "Add structures selected for displacement type, amplitude, and predicted response variation—not simply more random snapshots.", ORANGE),
            ("3", "Matched model comparisons", "Use the same initialization and split to isolate fixed physics, learned tensors, and direct Full-H training.", BLUE),
            ("4", "Denser vibration grid", "Resolve the long-wavelength polar limit after the tensor contribution is independently detectable.", GREEN),
        ]
        for i, (num, head, body, color) in enumerate(next_steps):
            x = 0.62 + (i % 2) * 4.48
            y = 1.64 + (i // 2) * 2.16
            box(s, x, y, 4.20, 1.78, fill=LIGHT, line=LINE, radius=True)
            box(s, x + 0.20, y + 0.21, 0.44, 0.44, fill=color, radius=True)
            text(s, num, x + 0.20, y + 0.27, 0.44, 0.19, size=11,
                 color=WHITE, bold=True, align=PP_ALIGN.CENTER)
            text(s, head, x + 0.80, y + 0.20, 3.16, 0.31, size=15.5,
                 color=NAVY, font="Cambria", bold=True)
            text(s, body, x + 0.80, y + 0.68, 3.12, 0.82, size=10.3,
                 color=GRAY)
        box(s, 0.84, 6.18, 8.32, 0.52, fill=NAVY, radius=True)
        text(s, "Scientific decision gate: claim a tensor benefit only when learned − fixed exceeds finite-difference uncertainty reproducibly.",
             1.08, 6.31, 7.84, 0.24, size=9.6, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER)
        footer(s, 14)
        notes(s, "Make the tensor effect measurable before scaling up", """
12:30–13:35. My first next step would be a displacement-sensitivity ladder,
because the current five-microångström step may make geometry-dependent tensor
changes invisible. Then I would add response-diverse labels and repeat matched
ablations with the same initialization. A denser q grid is valuable, but only
after the tensor contribution is measurable independently of numerical noise.
""")

        # 15 — summary and acknowledgments
        s = new_slide(prs)
        title(s, "Summary", "Physics-guided learning improves response—even before tensors help", 15)
        takeaways = [
            ("1", "I built a shared model for local electronic structure, effective charge, and screening.", BLUE),
            ("2", "The new response heads beat constant baselines on locked structures.", GREEN),
            ("3", "Direct Full-H wins on matrix values; both SR models win on atomic-motion response.", ORANGE),
            ("4", "Learned tensors have not changed EPC at the present displacement sensitivity.", TEAL),
        ]
        for i, (num, body, color) in enumerate(takeaways):
            y = 1.63 + i * 0.83
            box(s, 0.72, y, 0.39, 0.39, fill=color, radius=True)
            text(s, num, 0.72, y + 0.06, 0.39, 0.18, size=10,
                 color=WHITE, bold=True, align=PP_ALIGN.CENTER)
            text(s, body, 1.36, y - 0.01, 5.06, 0.52, size=12.3,
                 color=NAVY, bold=True, valign=MSO_ANCHOR.MIDDLE)
        box(s, 6.77, 1.62, 2.58, 3.92, fill=LIGHT, line=LINE, radius=True)
        text(s, "Acknowledgments", 7.03, 1.94, 2.06, 0.34, size=18,
             color=NAVY, font="Cambria", bold=True,
             align=PP_ALIGN.CENTER)
        ack = "Prof. Marco Benardi\n\nYao Luo, graduate mentor\n\nWAVE at Caltech"
        text(s, ack, 7.05, 2.61, 2.02, 2.25, size=11.5, color=GRAY,
             bold=True, align=PP_ALIGN.CENTER,
             valign=MSO_ANCHOR.MIDDLE)
        add_picture_fit(s, logo, 7.23, 4.86, 1.66, 0.48, border=None)
        box(s, 0.92, 5.52, 5.40, 0.71, fill=NAVY, radius=True)
        text(s, "The broader lesson: evaluate a scientific model on the physical response you care about—not only the values it fits.",
             1.18, 5.69, 4.88, 0.34, size=10.3, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER)
        text(s, "Thank you", 0.92, 6.48, 5.40, 0.34, size=18,
             color=TEAL, font="Cambria", bold=True,
             align=PP_ALIGN.CENTER)
        footer(s, 15)
        notes(s, "Physics-guided learning improves response—even before tensors help", """
13:35–14:35. In summary, I built the intended three-output architecture and
verified it on locked data. The response heads learn real signal. Direct
full-H remains best for matrix values, but target design around short-range
physics produces far better atomic-motion response. Learned tensors have not
yet improved EPC at this sensitivity. Thank Prof. Marco Benardi, my graduate
mentor Yao Luo, and WAVE at Caltech for their support.
""")

        # 16 — required final feedback slide
        s = new_slide(prs, WHITE)
        box(s, 0, 0, 0.20, H, fill=MINT)
        text(s, "Thank you", 0.72, 0.58, 4.38, 0.52, size=31,
             color=NAVY, font="Cambria", bold=True)
        text(s, "Questions?", 0.72, 1.30, 4.38, 0.42, size=22,
             color=TEAL, font="Cambria", bold=True)
        text(s, "Jeremiah Bailey", 0.72, 2.18, 4.22, 0.43, size=21,
             color=NAVY, font="Cambria", bold=True)
        text(s, "Please scan the QR code to leave\naudience feedback for this talk.",
             0.72, 2.93, 4.12, 0.86, size=15, color=GRAY, bold=True)
        box(s, 0.72, 4.27, 3.96, 0.62, fill=NAVY, radius=True)
        text(s, "Your feedback will be shared after the presentations.",
             0.94, 4.45, 3.52, 0.25, size=9.7, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER)
        add_picture_fit(s, logo, 0.93, 5.46, 2.10, 0.62, border=None)
        add_picture_fit(s, qr, 5.18, 0.72, 4.10, 5.95, border=LINE)
        text(s, "Audience feedback", 5.43, 6.75, 3.60, 0.27, size=12,
             color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        notes(s, "Questions and audience feedback", """
14:35–15:00. Thank the audience and invite questions. Leave this final slide
displayed throughout the 2–3 minute Q&A so the QR code remains available.
""")

        prs.save(OUTPUT)

    note_lines = [
        "# Jeremiah Bailey — SURF final presentation speaker notes", "",
        "Format: 4:3  ",
        "Target speaking time: 15 minutes  ",
        "Q&A: 2–3 minutes  ",
        "Audience: general Caltech audience", "",
    ]
    for index, (heading, body) in enumerate(theme.SLIDE_NOTES, 1):
        note_lines.extend([f"## {index}. {heading}", "", body, ""])
    NOTES_OUTPUT.write_text("\n".join(note_lines), encoding="utf-8")
    print(OUTPUT)
    print(NOTES_OUTPUT)


if __name__ == "__main__":
    build()
