#!/usr/bin/env python3
"""Build the August 17 learned-response results presentation."""
from __future__ import annotations

import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches

ROOT = Path(__file__).resolve().parents[2]
HERE = ROOT / "results" / "learned_response"

from workflows.analysis.build_project_update_presentation import (
    BLUE, GRAY, GREEN, INK, LIGHT, LINE, MID, MINT, NAVY, ORANGE,
    PALE_BLUE, PALE_ORANGE, RED, TEAL, WHITE, add_notes,
    add_picture_fit, box, line, metric, new_slide, rich_text, status_chip,
    text,
)
import workflows.analysis.build_project_update_presentation as theme


OUTPUT = HERE / "MACE-H-LR_Learned_Response_Results_2026-08-17.pptx"
NOTES_OUTPUT = HERE / "MACE-H-LR_Learned_Response_Results_2026-08-17_speaker_notes.md"
METRICS = HERE / "metrics.json"
W, H = 10.0, 5.625


def slide_title(slide, kicker, headline, number):
    text(slide, kicker.upper(), 0.6, 0.25, 8.7, 0.20, size=9.5,
         color=TEAL, bold=True)
    size = 23 if len(headline) > 50 else 27
    text(slide, headline, 0.6, 0.48, 8.8, 0.62, size=size, color=NAVY,
         font="Cambria", bold=True)
    line(slide, 0.6, 1.19, 9.4, 1.19, color=LINE, width=0.8)
    text(slide, str(number), 9.12, 0.28, 0.28, 0.16, size=8.5,
         color=GRAY, align=PP_ALIGN.RIGHT)


def slide_footer(slide, number):
    text(slide, "MACE-H-LR  ·  LEARNED RESPONSE RESULTS  ·  17 AUG 2026",
         0.6, 5.36, 5.8, 0.12, size=7.2, color=GRAY)
    text(slide, str(number), 9.08, 5.36, 0.32, 0.12, size=7.2,
         color=GRAY, align=PP_ALIGN.RIGHT)


def picture(slide, name, x, y, w, h):
    return add_picture_fit(slide, HERE / name, x, y, w, h)


def percent_reduction(value, baseline):
    return 100 * (1 - value / baseline)


def comparison_bar(slide, label, value, baseline, x, y, w, color, units):
    maximum = max(value, baseline)
    text(slide, label, x, y, w, 0.20, size=10, color=NAVY, bold=True)
    text(slide, "training-mean baseline", x, y + 0.34, 1.45, 0.16,
         size=8.2, color=GRAY)
    box(slide, x + 1.54, y + 0.34, w - 1.54, 0.16, fill=LINE, radius=True)
    box(slide, x + 1.54, y + 0.34,
        (w - 1.54) * baseline / maximum, 0.16, fill=GRAY, radius=True)
    text(slide, f"{baseline:.6f} {units}", x + 1.54, y + 0.54,
         w - 1.54, 0.16, size=8.2, color=GRAY, align=PP_ALIGN.RIGHT)
    text(slide, "learned head", x, y + 0.89, 1.45, 0.16,
         size=8.2, color=color, bold=True)
    box(slide, x + 1.54, y + 0.89, w - 1.54, 0.16, fill=LINE, radius=True)
    box(slide, x + 1.54, y + 0.89,
        (w - 1.54) * value / maximum, 0.16, fill=color, radius=True)
    text(slide, f"{value:.6f} {units}", x + 1.54, y + 1.09,
         w - 1.54, 0.16, size=8.2, color=color, bold=True,
         align=PP_ALIGN.RIGHT)


def build():
    data = json.loads(METRICS.read_text())
    h = data["hamiltonian"]
    epc = data["epc"]
    tensors = data["tensor_targets"]
    born = tensors["born"]
    epsilon = tensors["epsilon"]
    baseline = tensors["constant_baseline"]

    theme.SLIDE_NOTES.clear()
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    # 1 — title
    s = new_slide(prs, NAVY)
    box(s, 0, 0, 0.18, H, fill=MINT)
    status_chip(s, "Results update", 0.75, 0.62, 1.40, color=MINT,
                text_color=NAVY)
    text(s, "Learned response tensors", 0.75, 1.23, 8.55, 0.62,
         size=34, color=WHITE, font="Cambria", bold=True)
    text(s, "What Born-charge and dielectric heads change—and what they do not",
         0.75, 2.03, 8.25, 0.82, size=20.5, color=PALE_BLUE,
         font="Cambria", bold=True)
    rich_text(s, [
        ("Hamiltonian  ·  ", MID, False, 13),
        ("tensor generalization  ·  ", MINT, True, 13),
        ("Cartesian-AO EPC", MID, False, 13),
    ], 0.75, 3.15, 8.5, 0.32)
    text(s, "Four-way comparison: actual DFT · direct Full-H · LR-corrected SR · SR + predicted Z* + ε∞",
         0.75, 4.16, 8.45, 0.42, size=10.8, color=MID, bold=True)
    text(s, "Bernardi Group, Caltech  ·  August 17, 2026",
         0.75, 4.76, 8.5, 0.22, size=10, color=MID)
    add_notes(s, "Learned response tensors", """
0:00–0:35. This update closes the loop on the learned Born-charge and
dielectric architecture. I will show three levels of evidence: whether the
heads generalize on locked tensor labels, whether reconstructed full-H values
improve, and whether geometry-dependent tensors change Cartesian-AO EPC.
""")

    # 2 — executive result
    s = new_slide(prs)
    slide_title(s, "Executive result", "The heads learned the tensors; EPC is unchanged so far", 2)
    metric(s, f"{born['mae']:.4f} e", "LOCKED BORN MAE", 0.65, 1.46, 2.02,
           color=ORANGE, fill=PALE_ORANGE)
    metric(s, f"{epsilon['mae']:.6f}", "LOCKED ε∞ MAE", 2.83, 1.46, 2.02,
           color=GREEN, fill="E4F4EE")
    metric(s, f"{1e3*h['tensor']['mae']:.3f} meV", "RECONSTRUCTED FULL-H MAE",
           5.01, 1.46, 2.08, color=BLUE, fill="E7F2F8")
    metric(s, f"{100*epc['tensor']['relative_l2']:.2f}%", "PREDICTED-TENSOR EPC REL. L2",
           7.25, 1.46, 2.10, color=TEAL)
    conclusions = [
        ("LEARNED", "Born and ε∞ beat training-mean baselines on five locked structures.", GREEN),
        ("PRESERVED", "Ten consecutive validation passes preserved SR-H accuracy.", BLUE),
        ("NEUTRAL", "Predicted tensors change full-H values and EPC only marginally here.", ORANGE),
        ("STILL TRUE", "Direct Full-H is better on H values but much worse on EPC response.", RED),
    ]
    for i, (tag, body, color) in enumerate(conclusions):
        y = 2.58 + i * 0.57
        status_chip(s, tag, 0.70, y, 1.05, color=color)
        text(s, body, 1.95, y + 0.02, 7.20, 0.24, size=11.5,
             color=INK, bold=(i == 2))
    slide_footer(s, 2)
    add_notes(s, "The heads learned the tensors; EPC is unchanged so far", """
0:35–1:25. The headline is nuanced. The tensor heads are not trivial constant
predictors: both beat the locked baselines. Training preserved the short-range
Hamiltonian. But at the current five-microångström EPC displacement, replacing
fixed tensors with predicted geometry-dependent tensors produces essentially
the same EPC. The direct full-H model remains best for H values and worst for
the derivative observable.
""")

    # 3 — architecture and protocol
    s = new_slide(prs)
    slide_title(s, "Model + protocol", "A shared encoder now predicts HSR, Z*, and ε∞", 3)
    blocks = [
        (0.60, 1.83, 1.25, "Structure", "R + lattice", LIGHT, NAVY),
        (2.13, 1.65, 1.42, "MACE-H", "shared encoder", NAVY, WHITE),
        (4.01, 1.28, 1.46, "HSR", "edge head", "E7F2F8", BLUE),
        (4.01, 2.20, 1.46, "Z*", "atom head + ASR", PALE_ORANGE, ORANGE),
        (4.01, 3.12, 1.46, "ε∞", "global SPD head", "E4F4EE", GREEN),
        (6.05, 1.88, 1.48, "Analytic HLR", "predicted tensors", "E4F4EE", GREEN),
        (8.11, 1.88, 1.28, "Hfull", "HSR + HLR", NAVY, WHITE),
    ]
    for x, y, w, head, body, fill, color in blocks:
        box(s, x, y, w, 0.92, fill=fill, line=LINE, radius=True)
        text(s, head, x + 0.08, y + 0.15, w - 0.16, 0.25, size=15,
             color=color, font="Cambria", bold=True, align=PP_ALIGN.CENTER)
        text(s, body, x + 0.08, y + 0.54, w - 0.16, 0.18, size=8.5,
             color=MID if fill == NAVY else GRAY, align=PP_ALIGN.CENTER)
    line(s, 1.88, 2.29, 2.07, 2.29, color=TEAL, width=2, arrow=True)
    for y in (1.74, 2.66, 3.58):
        line(s, 3.59, 2.13, 3.94, y, color=TEAL, width=1.6, arrow=True)
    line(s, 5.51, 2.66, 5.99, 2.34, color=TEAL, width=1.8, arrow=True)
    line(s, 5.51, 3.58, 5.99, 2.50, color=TEAL, width=1.8, arrow=True)
    line(s, 5.51, 1.74, 8.03, 2.12, color=TEAL, width=1.8, arrow=True)
    line(s, 7.57, 2.34, 8.03, 2.34, color=TEAL, width=2, arrow=True)
    stages = [
        ("1", "Head-only", "freeze H + encoder", GREEN),
        ("2", "Partial", "unfreeze final block", BLUE),
        ("3", "Locked test", "open once after selection", ORANGE),
        ("4", "Two EPC modes", "frozen vs geometry-dependent", TEAL),
    ]
    for i, (num, head, body, color) in enumerate(stages):
        x = 0.66 + 2.18 * i
        box(s, x, 4.36, 1.92, 0.62, fill=LIGHT, line=LINE, radius=True)
        status_chip(s, num, x + 0.10, 4.52, 0.32, color=color)
        text(s, head, x + 0.51, 4.45, 1.25, 0.17, size=9.5,
             color=NAVY, bold=True)
        text(s, body, x + 0.51, 4.69, 1.25, 0.16, size=7.7, color=GRAY)
    slide_footer(s, 3)
    add_notes(s, "A shared encoder now predicts HSR, Z*, and ε∞", """
1:25–2:10. The architectural change is explicit: one equivariant encoder
branches to the short-range Hamiltonian, full atom-resolved Born tensors, and
a symmetric positive-definite dielectric tensor. Acoustic sum rule and SPD
constraints are built in. Training progressed from head-only to a limited
final-block unfreeze, and the locked test was opened only after the model and
learning rates were frozen.
""")

    # 4 — training figure
    s = new_slide(prs)
    slide_title(s, "Optimization", "Partial fine-tuning passed every gate for ten epochs", 4)
    picture(s, "01_training_and_multitask_gates.png", 0.48, 1.34, 9.04, 3.58)
    box(s, 0.78, 4.94, 8.44, 0.25, fill=NAVY, radius=True)
    text(s, "Selected epoch 9: H MSE 2.857×10⁻⁶ · Born MAE 1.799×10⁻³ e · ε∞ MAE 3.386×10⁻⁴",
         0.96, 5.00, 8.08, 0.14, size=8.5, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER)
    slide_footer(s, 4)
    add_notes(s, "Partial fine-tuning passed every gate for ten epochs", """
2:10–2:55. The left panel gives context from the original Hamiltonian runs.
The center panel shows that partial unfreezing did not push SR validation MSE
through the three-times-ten-to-the-minus-six ceiling. On the right, both tensor
tasks remain below their stricter zero-residual baselines for all ten epochs.
The stop is a consecutive-gate criterion, not selection from one lucky epoch.
""")

    # 5 — tensor test
    s = new_slide(prs)
    slide_title(s, "Locked tensor test", "The heads generalize beyond a constant predictor", 5)
    comparison_bar(s, "Born tensor", born["mae"], baseline["born_mae"],
                   0.72, 1.53, 4.10, ORANGE, "e")
    comparison_bar(s, "Electronic dielectric tensor", epsilon["mae"],
                   baseline["epsilon_mae"], 5.18, 1.53, 4.10, GREEN, "")
    metric(s, f"{percent_reduction(born['mae'], baseline['born_mae']):.1f}%",
           "BORN MAE REDUCTION", 0.83, 3.31, 1.82, color=ORANGE,
           fill=PALE_ORANGE)
    metric(s, f"{percent_reduction(epsilon['mae'], baseline['epsilon_mae']):.1f}%",
           "ε∞ MAE REDUCTION", 2.86, 3.31, 1.82, color=GREEN,
           fill="E4F4EE")
    metric(s, "1.88×10⁻⁵ e", "MAX BORN ASR RESIDUAL", 5.05, 3.31, 1.98,
           color=BLUE, fill="E7F2F8")
    metric(s, "3.326–3.333", "PREDICTED ε∞ EIGENVALUES", 7.25, 3.31, 2.02,
           color=TEAL)
    box(s, 0.82, 4.42, 8.36, 0.52, fill=LIGHT, line=LINE, radius=True)
    text(s, "Evidence of geometry-dependent learning—not production generalization: 10 train · 2 validation · 5 locked test structures.",
         1.02, 4.58, 7.96, 0.20, size=10.2, color=NAVY, bold=True,
         align=PP_ALIGN.CENTER)
    slide_footer(s, 5)
    add_notes(s, "The heads generalize beyond a constant predictor", """
2:55–3:45. On the locked structures, Born MAE is 0.002369 electron versus a
0.003201 training-mean baseline, a 26 percent reduction. Dielectric MAE falls
by 73 percent. The physical outputs remain constrained: Born satisfies ASR to
float32 accumulation precision, and all dielectric eigenvalues are positive.
This is evidence that the heads learned signal, but only five locked examples
are not enough for a production claim.
""")

    # 6 — Hamiltonian
    s = new_slide(prs)
    slide_title(s, "Hamiltonian values", "Predicted tensors barely change full-H accuracy", 6)
    picture(s, "02_locked_hamiltonian_error.png", 0.50, 1.35, 7.10, 3.65)
    metric(s, f"{1e3*h['full']['mae']:.3f}", "DIRECT FULL-H MAE (meV)",
           7.78, 1.56, 1.63, color=ORANGE, fill=PALE_ORANGE)
    metric(s, f"{1e3*h['sr']['mae']:.3f}", "FIXED-LR SR MAE (meV)",
           7.78, 2.57, 1.63, color=BLUE, fill="E7F2F8")
    metric(s, f"{1e3*h['tensor']['mae']:.3f}", "PREDICTED-LR SR MAE (meV)",
           7.78, 3.58, 1.63, color=GREEN, fill="E4F4EE")
    text(s, "Full-H reconstruction identity: exact\nBlock coverage: 123,716 / 123,716",
         7.79, 4.60, 1.60, 0.42, size=8.5, color=GRAY,
         align=PP_ALIGN.CENTER)
    slide_footer(s, 6)
    add_notes(s, "Predicted tensors barely change full-H accuracy", """
3:45–4:30. On the same five locked structures, direct full-H remains best on
matrix elements at 0.417 meV MAE. The old fixed-LR SR reconstruction is 0.487,
and the predicted-tensor reconstruction is 0.485. So the new tensors make a
small improvement to the SR reconstruction, but they do not close the value
gap. The reconstruction is nevertheless complete and exact on every truth
block.
""")

    # 7 — EPC overall
    s = new_slide(prs)
    slide_title(s, "EPC overall", "SR still wins the response test; learned tensors are neutral", 7)
    picture(s, "04_epc_overall_error.png", 0.50, 1.35, 7.12, 3.66)
    box(s, 7.82, 1.52, 1.54, 1.02, fill=NAVY, radius=True)
    text(s, "24.36%", 7.94, 1.71, 1.30, 0.28, size=21, color=PALE_BLUE,
         font="Cambria", bold=True, align=PP_ALIGN.CENTER)
    text(s, "fixed-LR SR", 7.94, 2.14, 1.30, 0.16, size=8.5, color=MID,
         bold=True, align=PP_ALIGN.CENTER)
    box(s, 7.82, 2.75, 1.54, 1.02, fill="E4F4EE", line=LINE, radius=True)
    text(s, "24.43%", 7.94, 2.94, 1.30, 0.28, size=21, color=GREEN,
         font="Cambria", bold=True, align=PP_ALIGN.CENTER)
    text(s, "predicted tensors", 7.94, 3.37, 1.30, 0.16, size=8.5,
         color=GRAY, bold=True, align=PP_ALIGN.CENTER)
    box(s, 7.82, 4.00, 1.54, 0.80, fill=PALE_ORANGE, line=LINE, radius=True)
    text(s, "91.93%", 7.94, 4.15, 1.30, 0.24, size=18, color=ORANGE,
         font="Cambria", bold=True, align=PP_ALIGN.CENTER)
    text(s, "direct Full-H", 7.94, 4.51, 1.30, 0.14, size=8, color=GRAY,
         align=PP_ALIGN.CENTER)
    slide_footer(s, 7)
    add_notes(s, "SR still wins the response test; learned tensors are neutral", """
4:30–5:15. The central EPC result remains stable. Direct full-H has 91.93
percent relative L2 error. Both SR pipelines are around 24.4 percent. But the
new geometry-dependent tensor pipeline is slightly worse, not better, than
fixed-reference tensors by 0.077 percentage points. At this sensitivity, the
new heads have neither helped nor damaged the EPC conclusion.
""")

    # 8 — q result
    s = new_slide(prs)
    slide_title(s, "EPC by q", "The two SR curves overlap across the full 2×2×2 grid", 8)
    picture(s, "05_epc_q_resolved.png", 0.47, 1.32, 9.08, 3.78)
    box(s, 0.78, 4.96, 8.44, 0.25, fill=NAVY, radius=True)
    text(s, "Direct Full-H overestimates finite-q strength; predicted and fixed tensors are visually coincident.",
         0.98, 5.02, 8.04, 0.14, size=8.5, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER)
    slide_footer(s, 8)
    add_notes(s, "The two SR curves overlap across the full 2×2×2 grid", """
5:15–6:00. The q-resolved view shows where the separation occurs. Direct
full-H greatly overestimates EPC strength at finite q. The old fixed-tensor SR
and new predicted-tensor SR curves overlap at every sampled q. Different line
styles and markers are used because the numerical differences are otherwise
too small to see.
""")

    # 9 — parity and distributions
    s = new_slide(prs)
    slide_title(s, "EPC structure", "Predicted tensors preserve the SR response distribution", 9)
    picture(s, "06_epc_dft_parity.png", 0.50, 1.34, 6.08, 3.66)
    picture(s, "07_epc_magnitude_distribution.png", 6.78, 1.34, 2.70, 2.73)
    box(s, 6.82, 4.22, 2.62, 0.74, fill=LIGHT, line=LINE, radius=True)
    text(s, "The green and blue magnitude CDFs overlap. Both miss the same low-magnitude DFT structure, while direct Full-H broadens the response.",
         7.00, 4.34, 2.26, 0.48, size=8.8, color=INK,
         align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    slide_footer(s, 9)
    add_notes(s, "Predicted tensors preserve the SR response distribution", """
6:00–6:45. The parity view confirms the same hierarchy: direct full-H has a
much broader cloud around the DFT identity line, while both SR variants are
tighter and nearly indistinguishable. The magnitude distributions show that
the tensor change does not alter the current response manifold; both SR runs
retain the same systematic differences from DFT.
""")

    # 10 — interpretation
    s = new_slide(prs)
    slide_title(s, "Interpretation", "Three statements are supported—and three are not yet", 10)
    box(s, 0.62, 1.45, 4.25, 3.55, fill="E4F4EE", line=LINE, radius=True)
    status_chip(s, "Supported", 0.86, 1.69, 1.00, color=GREEN)
    supported = [
        "Tensor heads learn held-out response variation.",
        "The composed SR + analytic-LR contract works exactly.",
        "SR training produces a better EPC response than direct full-H on this grid.",
    ]
    for i, value in enumerate(supported):
        y = 2.22 + i * 0.78
        box(s, 0.88, y, 0.28, 0.28, fill=GREEN, radius=True)
        text(s, "✓", 0.88, y + 0.03, 0.28, 0.15, size=9.5, color=WHITE,
             bold=True, align=PP_ALIGN.CENTER)
        text(s, value, 1.34, y - 0.01, 3.12, 0.48, size=11,
             color=NAVY, bold=True)
    box(s, 5.13, 1.45, 4.25, 3.55, fill=PALE_ORANGE, line=LINE, radius=True)
    status_chip(s, "Not established", 5.37, 1.69, 1.33, color=ORANGE)
    open_questions = [
        "Geometry-dependent tensors improve EPC.",
        "Ten tensor-training structures generalize broadly.",
        "The present 2×2×2 grid resolves the small-q polar limit.",
    ]
    for i, value in enumerate(open_questions):
        y = 2.22 + i * 0.78
        box(s, 5.39, y, 0.28, 0.28, fill=ORANGE, radius=True)
        text(s, "?", 5.39, y + 0.03, 0.28, 0.15, size=9.5, color=WHITE,
             bold=True, align=PP_ALIGN.CENTER)
        text(s, value, 5.85, y - 0.01, 3.12, 0.48, size=11,
             color=NAVY, bold=True)
    slide_footer(s, 10)
    add_notes(s, "Three statements are supported—and three are not yet", """
6:45–7:35. The clean conclusion is not that tensor learning failed. The heads
learned held-out structure dependence and the composed physics path works. The
result we do not have is an EPC improvement from geometry-dependent tensors.
We also cannot generalize from ten tensor-training structures, and a 2×2×2 q
grid is not a small-q convergence study. These boundaries matter for how we
describe the result.
""")

    # 11 — next experiments
    s = new_slide(prs)
    slide_title(s, "Next experiments", "Increase sensitivity before increasing model complexity", 11)
    steps = [
        ("1", "Finite-difference ladder", "Repeat EPC at 5×10⁻⁶, 10⁻⁴, 10⁻³, and 5×10⁻³ Å; separate numerical invisibility from physical neutrality.", TEAL),
        ("2", "Tensor-label coverage", "Add displacement families and amplitudes where predicted Z* and ε∞ vary most; keep train/validation families separated.", ORANGE),
        ("3", "Matched ablations", "Same initialization and seed: SR-only, fixed tensors, learned frozen tensors, learned geometry-dependent tensors, direct Full-H.", BLUE),
        ("4", "q-grid sensitivity", "Move beyond 2×2×2 only after derivative and tensor effects are measurable; target the finite-q failure directly.", GREEN),
    ]
    for i, (num, head, body, color) in enumerate(steps):
        x = 0.62 + (i % 2) * 4.48
        y = 1.45 + (i // 2) * 1.75
        box(s, x, y, 4.20, 1.43, fill=LIGHT, line=LINE, radius=True)
        box(s, x + 0.18, y + 0.18, 0.42, 0.42, fill=color, radius=True)
        text(s, num, x + 0.18, y + 0.23, 0.42, 0.20, size=11,
             color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        text(s, head, x + 0.76, y + 0.18, 3.18, 0.27, size=14.5,
             color=NAVY, font="Cambria", bold=True)
        text(s, body, x + 0.76, y + 0.57, 3.16, 0.65, size=9.5,
             color=GRAY)
    box(s, 0.82, 4.96, 8.36, 0.24, fill=NAVY, radius=True)
    text(s, "Decision gate: do not claim learned-tensor EPC benefit until C − B exceeds finite-difference uncertainty reproducibly.",
         1.00, 5.02, 8.00, 0.13, size=8.2, color=WHITE, bold=True,
         align=PP_ALIGN.CENTER)
    slide_footer(s, 11)
    add_notes(s, "Increase sensitivity before increasing model complexity", """
7:35–8:45. The immediate priority is a sensitivity experiment, not a larger
network. The geometry-dependent correction may be invisible because the EPC
displacement is only five microångströms. I would run a controlled displacement
ladder, then add tensor labels chosen for response variation, and use matched
initializations for the causal ablation. Dense q should come only after the
effect is resolvable above finite-difference uncertainty.
""")

    # 12 — takeaways
    s = new_slide(prs, NAVY)
    box(s, 0, 0, 0.18, H, fill=MINT)
    text(s, "What we see so far", 0.75, 0.58, 8.5, 0.50, size=30,
         color=WHITE, font="Cambria", bold=True)
    takeaways = [
        ("1", "The new Born and dielectric heads pass locked baselines and physical constraints.", GREEN),
        ("2", "Reconstructed H is complete and stable, but direct Full-H remains better on matrix-element values.", ORANGE),
        ("3", "Both SR pipelines remain far better than direct Full-H for Cartesian-AO EPC.", BLUE),
        ("4", "At 5×10⁻⁶ Å, predicted geometry-dependent tensors are numerically indistinguishable from fixed tensors.", MINT),
        ("5", "The next result should be a controlled sensitivity and data-coverage study—not a stronger claim.", WHITE),
    ]
    for i, (num, body, color) in enumerate(takeaways):
        y = 1.40 + i * 0.70
        box(s, 0.78, y, 0.38, 0.38, fill=color if color != WHITE else MID,
            radius=True)
        text(s, num, 0.78, y + 0.06, 0.38, 0.18, size=10.5,
             color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        text(s, body, 1.40, y - 0.01, 7.75, 0.43, size=13,
             color=WHITE if i == 4 else PALE_BLUE, bold=(i == 4),
             valign=MSO_ANCHOR.MIDDLE)
    text(s, "Questions · which sensitivity test should define the next compute campaign?",
         0.75, 5.04, 8.45, 0.23, size=11.5, color=MINT, bold=True)
    add_notes(s, "What we see so far", """
8:45–9:30. The tensor-head implementation is successful as a prototype. It
learns held-out tensor variation, preserves Hamiltonian stability, and closes
the intended composed architecture. The current scientific outcome is that it
does not yet change EPC. The next useful result is to establish whether that
neutrality is physical, data-limited, or hidden below the finite-difference
sensitivity. Then invite discussion on the next compute campaign.
""")

    prs.save(OUTPUT)
    lines = [
        "# Learned response tensors — speaker notes", "",
        "Presentation date: August 17, 2026  ",
        "Target length: 9–10 minutes", "",
    ]
    for index, (heading, notes) in enumerate(theme.SLIDE_NOTES, 1):
        lines.extend([f"## {index}. {heading}", "", notes, ""])
    NOTES_OUTPUT.write_text("\n".join(lines), encoding="utf-8")
    print(OUTPUT)
    print(NOTES_OUTPUT)


if __name__ == "__main__":
    build()
