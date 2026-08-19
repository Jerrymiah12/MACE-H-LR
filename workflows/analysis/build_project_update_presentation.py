#!/usr/bin/env python3
"""Build the August 13, 2026 MACE-H-LR project-update presentation.

The deck deliberately uses the evidence already collected in ``plots/`` and
keeps the visual language of ``MACE-H-LR_10min_template.pptx``.  Speaker notes
are embedded in the PowerPoint and also exported as Markdown.
"""

from __future__ import annotations

import tempfile
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
HERE = ROOT / "results" / "project_update"
TEMPLATE = Path(__file__).resolve().parent / "assets" / "MACE-H-LR_10min_template.pptx"
OUTPUT = HERE / "MACE-H-LR_Project_Update_2026-08-13.pptx"
NOTES_OUTPUT = HERE / "MACE-H-LR_Project_Update_2026-08-13_speaker_notes.md"

NAVY = "12224A"
TEAL = "1C7293"
MINT = "02C39A"
BLUE = "177FB5"
PALE_BLUE = "DCE6F5"
ORANGE = "D96000"
PALE_ORANGE = "F8E3D2"
RED = "B53A3A"
GREEN = "2D8A66"
INK = "243447"
GRAY = "5A6B7B"
MID = "B9C6DE"
LIGHT = "F4F7FA"
LINE = "D8E0EA"
WHITE = "FFFFFF"

W = 10.0
H = 5.625
SLIDE_NOTES: list[tuple[str, str]] = []


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def remove_all_slides(prs: Presentation) -> None:
    slide_ids = prs.slides._sldIdLst  # noqa: SLF001 - no public delete API
    for slide_id in list(slide_ids):
        slide_ids.remove(slide_id)


def set_bg(slide, color: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def box(slide, x, y, w, h, fill=WHITE, line=None, radius=False, transparency=0):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    shp.fill.transparency = transparency
    if line:
        shp.line.color.rgb = rgb(line)
        shp.line.width = Pt(0.8)
    else:
        shp.line.fill.background()
    if radius:
        try:
            shp.adjustments[0] = 0.08
        except (IndexError, ValueError):
            pass
    return shp


def line(slide, x1, y1, x2, y2, color=LINE, width=1.5, arrow=False):
    shp = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    shp.line.color.rgb = rgb(color)
    shp.line.width = Pt(width)
    if arrow:
        # python-pptx exposes no stable public arrowhead API; use a small
        # triangle at the end so the diagram remains portable.
        tri = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            Inches(x2 - 0.08), Inches(y2 - 0.065), Inches(0.13), Inches(0.13),
        )
        tri.rotation = 90
        tri.fill.solid()
        tri.fill.fore_color.rgb = rgb(color)
        tri.line.fill.background()
    return shp


def text(
    slide,
    value,
    x,
    y,
    w,
    h,
    *,
    size=14,
    color=INK,
    font="Calibri",
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.0,
    italic=False,
    fit=False,
    line_spacing=1.0,
):
    shp = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shp.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    r = p.add_run()
    r.text = value
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = rgb(color)
    if fit:
        tf.fit_text(font_family=font, max_size=size)
    return shp


def rich_text(slide, runs, x, y, w, h, *, size=14, align=PP_ALIGN.LEFT,
              valign=MSO_ANCHOR.TOP, margin=0.0, font="Calibri"):
    shp = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shp.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    for value, color, bold, run_size in runs:
        r = p.add_run()
        r.text = value
        r.font.name = font
        r.font.size = Pt(run_size or size)
        r.font.bold = bold
        r.font.color.rgb = rgb(color)
    return shp


def title(slide, kicker: str, headline: str, number: int | None = None) -> None:
    text(slide, kicker.upper(), 0.6, 0.27, 8.8, 0.22, size=10.5,
         color=TEAL, bold=True)
    headline_size = 24 if len(headline) > 47 else 28
    headline_height = 0.76 if len(headline) > 47 else 0.62
    text(slide, headline, 0.6, 0.52, 8.8, headline_height,
         size=headline_size, color=NAVY,
         font="Cambria", bold=True)
    line(slide, 0.6, 1.31, 9.4, 1.31, color=LINE, width=0.8)
    if number is not None:
        text(slide, str(number), 9.18, 0.31, 0.22, 0.18, size=9,
             color=GRAY, align=PP_ALIGN.RIGHT)


def footer(slide, number: int) -> None:
    text(slide, "MACE-H-LR  ·  PROJECT UPDATE  ·  13 AUG 2026",
         0.6, 5.34, 5.4, 0.14, size=7.5, color=GRAY)
    text(slide, str(number), 9.05, 5.34, 0.35, 0.14, size=7.5,
         color=GRAY, align=PP_ALIGN.RIGHT)


def status_chip(slide, label, x, y, w, color=GREEN, text_color=WHITE):
    box(slide, x, y, w, 0.28, fill=color, radius=True)
    text(slide, label.upper(), x, y + 0.01, w, 0.20, size=8.2,
         color=text_color, bold=True, align=PP_ALIGN.CENTER,
         valign=MSO_ANCHOR.MIDDLE)


def metric(slide, value, label, x, y, w, color=TEAL, fill=LIGHT):
    box(slide, x, y, w, 0.78, fill=fill, line=LINE, radius=True)
    text(slide, value, x + 0.12, y + 0.10, w - 0.24, 0.34, size=21,
         color=color, bold=True, font="Cambria")
    text(slide, label, x + 0.12, y + 0.49, w - 0.24, 0.18, size=8.6,
         color=GRAY, bold=True)


def add_picture_fit(slide, path: Path, x, y, w, h, *, border=LINE):
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(w / iw, h / ih)
    pw, ph = iw * scale, ih * scale
    px, py = x + (w - pw) / 2, y + (h - ph) / 2
    if border:
        box(slide, x, y, w, h, fill=WHITE, line=border, radius=True)
    return slide.shapes.add_picture(str(path), Inches(px), Inches(py),
                                    Inches(pw), Inches(ph))


def crop_asset(source: Path, fractions, temp_dir: Path, name: str) -> Path:
    left, top, right, bottom = fractions
    with Image.open(source) as im:
        iw, ih = im.size
        out = im.crop((int(iw * left), int(ih * top),
                       int(iw * right), int(ih * bottom)))
        path = temp_dir / f"{name}.png"
        out.save(path)
    return path


def add_notes(slide, slide_title: str, notes: str) -> None:
    frame = slide.notes_slide.notes_text_frame
    frame.text = notes.strip()
    SLIDE_NOTES.append((slide_title, notes.strip()))


def new_slide(prs, bg=WHITE):
    layout = next(
        (candidate for candidate in prs.slide_layouts
         if candidate.name.lower() == "blank"),
        prs.slide_layouts[-1],
    )
    slide = prs.slides.add_slide(layout)
    set_bg(slide, bg)
    return slide


def build() -> None:
    # Build a clean OOXML package.  The palette and typography intentionally
    # follow the old deck, while avoiding stale slide parts from the template.
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    with tempfile.TemporaryDirectory(prefix="maceh_presentation_") as tmp:
        tmp_dir = Path(tmp)

        # 1 — title
        s = new_slide(prs, NAVY)
        box(s, 0.0, 0.0, 0.18, H, fill=MINT)
        status_chip(s, "Project update", 0.75, 0.68, 1.45, color=MINT,
                    text_color=NAVY)
        text(s, "MACE-H-LR", 0.75, 1.27, 8.4, 0.66, size=37,
             color=WHITE, font="Cambria", bold=True)
        text(s, "From validated machinery to a physics-correct full-H model",
             0.75, 2.06, 8.25, 0.82, size=22, color=PALE_BLUE,
             font="Cambria", bold=True)
        rich_text(s, [
            ("Gold machinery  ✓   ", MINT, True, 14),
            ("Cartesian EPC  ✓   ", MINT, True, 14),
            ("Learned Z* + ε∞ heads  →", "F5B65B", True, 14),
        ], 0.75, 3.23, 8.5, 0.38)
        text(s, "SURF 2026  ·  Bernardi Group, Caltech  ·  August 13, 2026",
             0.75, 4.35, 8.5, 0.30, size=12, color=MID, bold=True)
        text(s, "Handoff milestone: August 20", 0.75, 4.76, 8.5, 0.25,
             size=10.5, color=MID)
        add_notes(s, "MACE-H-LR", """
0:00–0:30. The project has moved through three layers: I first tested the
model machinery on bulk gold, then built and validated Cartesian
electron–phonon coupling, and the MgO results exposed an architecture issue.
The rest of the week is about correcting that issue with learned Born-charge
and dielectric heads and leaving a reproducible handoff by August 20.
""")

        # 2 — arc
        s = new_slide(prs)
        title(s, "Project arc", "Three stages changed the question", 2)
        stages = [
            ("1", "Bulk Au", "Machinery test", "Model + data path exercised", GREEN),
            ("2", "MgO Cartesian EPC", "Physics test", "301,056 complex AO components", BLUE),
            ("3", "Composed full-H", "Correction in progress", "Learn Z* and ε∞; reconstruct H", ORANGE),
        ]
        xs = [0.72, 3.68, 6.64]
        for idx, ((num, name, role, body, color), x) in enumerate(zip(stages, xs)):
            box(s, x, 1.63, 2.65, 2.66, fill=LIGHT, line=LINE, radius=True)
            box(s, x + 0.18, 1.86, 0.42, 0.42, fill=color, radius=True)
            text(s, num, x + 0.18, 1.91, 0.42, 0.24, size=12, color=WHITE,
                 bold=True, align=PP_ALIGN.CENTER)
            text(s, name, x + 0.18, 2.42, 2.25, 0.48, size=18.5, color=NAVY,
                 font="Cambria", bold=True)
            text(s, role.upper(), x + 0.18, 2.93, 2.25, 0.18, size=8.5,
                 color=color, bold=True)
            text(s, body, x + 0.18, 3.30, 2.25, 0.62, size=12,
                 color=GRAY, valign=MSO_ANCHOR.MIDDLE)
            if idx < 2:
                line(s, x + 2.67, 2.96, x + 2.91, 2.96, color=TEAL,
                     width=2.5, arrow=True)
        box(s, 0.72, 4.55, 8.57, 0.47, fill=NAVY, radius=True)
        text(s, "Outcome so far: the pipeline is real; the bottleneck is now the model objective.",
             0.93, 4.66, 8.15, 0.22, size=13, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER)
        footer(s, 2)
        add_notes(s, "Three stages changed the question", """
0:30–1:10. Stage one was the gold machinery test. Stage two was the first
end-to-end Cartesian-AO EPC calculation for MgO. That calculation did not
just produce a metric—it changed the question. The remaining problem is not
whether the plumbing runs; it is whether the training objective represents
the intended composed long-range model.
""")

        # 3 — target architecture / physics
        s = new_slide(prs)
        title(s, "Physics target", "The intended model is composed—not a direct full-H fit", 3)
        text(s, "Hfull(R) = HSR(R) + HLR[R; Z*(R), ε∞(R)]", 0.85, 1.47, 8.3, 0.46,
             size=24, color=NAVY, font="Cambria", bold=True,
             align=PP_ALIGN.CENTER)
        y = 2.32
        items = [
            (0.65, 1.72, "Structure", "positions + lattice", LIGHT, NAVY),
            (2.73, 1.58, "MACE-H", "shared equivariant backbone", LIGHT, NAVY),
            (4.66, 1.40, "HSR head", "local Hamiltonian", "E7F2F8", BLUE),
            (6.34, 1.37, "Z* + ε∞", "tensor heads", PALE_ORANGE, ORANGE),
            (8.00, 1.34, "HLR", "analytic physics", "E4F4EE", GREEN),
        ]
        for x, w, head, body, fill, color in items:
            box(s, x, y, w, 1.12, fill=fill, line=LINE, radius=True)
            text(s, head, x + 0.10, y + 0.18, w - 0.20, 0.28, size=16,
                 color=color, font="Cambria", bold=True, align=PP_ALIGN.CENTER)
            text(s, body, x + 0.10, y + 0.61, w - 0.20, 0.27, size=9.5,
                 color=GRAY, align=PP_ALIGN.CENTER)
        line(s, 2.37, 2.88, 2.67, 2.88, color=TEAL, width=2.2, arrow=True)
        line(s, 4.33, 2.65, 4.58, 2.65, color=TEAL, width=2.2, arrow=True)
        line(s, 4.33, 3.12, 6.26, 3.12, color=TEAL, width=2.2, arrow=True)
        line(s, 7.73, 2.88, 7.94, 2.88, color=TEAL, width=2.2, arrow=True)
        box(s, 0.85, 3.91, 8.30, 0.73, fill=NAVY, radius=True)
        rich_text(s, [
            ("Differentiate the reconstruction:  ", WHITE, False, 13),
            ("∂Hfull/∂τ", MINT, True, 15),
            ("  →  Cartesian EPC", WHITE, True, 13),
        ], 1.10, 4.12, 7.80, 0.30, align=PP_ALIGN.CENTER)
        text(s, "The network learns chemistry; the analytic term carries the nonlocal polar physics.",
             0.85, 4.82, 8.30, 0.30, size=12, color=GRAY,
             align=PP_ALIGN.CENTER, italic=True)
        footer(s, 3)
        add_notes(s, "The intended model is composed—not a direct full-H fit", """
1:10–2:00. This is the target architecture. The short-range Hamiltonian and
the two response tensors depend on structure. The analytic long-range term is
then reconstructed from the predicted Born charges and electronic dielectric
tensor. Cartesian EPC is the derivative of this entire composition. A direct
full-H network is a useful baseline, but it is not this model.
""")

        # 4 — training values looked fine
        s = new_slide(prs)
        title(s, "Initial training", "Hamiltonian metrics looked excellent—and hid the response problem", 4)
        loss = crop_asset(HERE / "figure_1_training_validation_loss.png",
                          (0.00, 0.03, 1.00, 0.94), tmp_dir, "loss")
        parity = crop_asset(HERE / "figure_2_test_parity.png",
                            (0.00, 0.04, 1.00, 0.94), tmp_dir, "parity")
        add_picture_fit(s, loss, 0.55, 1.43, 4.35, 2.18)
        add_picture_fit(s, parity, 5.10, 1.43, 4.35, 2.18)
        metric(s, "0.294 meV", "DIRECT FULL-H HELD-OUT MAE", 0.74, 3.85, 2.55,
               color=ORANGE, fill=PALE_ORANGE)
        metric(s, "0.389 meV", "LR-CORRECTED SR HELD-OUT MAE", 3.51, 3.85, 2.70,
               color=BLUE, fill="E7F2F8")
        metric(s, "≈ 0.99999", "R² FOR BOTH MODELS", 6.43, 3.85, 2.82,
               color=TEAL)
        box(s, 0.74, 4.84, 8.51, 0.30, fill=LIGHT, radius=True)
        text(s, "A low value loss does not guarantee an accurate displacement derivative.",
             0.90, 4.91, 8.20, 0.16, size=10.5, color=NAVY, bold=True,
             align=PP_ALIGN.CENTER)
        footer(s, 4)
        add_notes(s, "Hamiltonian metrics looked excellent—and hid the response problem", """
2:00–2:45. Both runs converged cleanly. On 37 held-out snapshots, the direct
full-H model even had the lower matrix-element MAE, and both parity plots look
nearly perfect. If I had stopped at Hamiltonian values, I would have concluded
that the direct full-H baseline was better. EPC tests the local slope with
respect to atomic displacement, which is a stricter observable.
""")

        # 5 — main EPC result
        s = new_slide(prs)
        title(s, "Response test", "Cartesian EPC exposed a finite-q failure in direct full-H", 5)
        qplot = crop_asset(HERE / "epc_01_q_resolved_comparison.png",
                           (0.01, 0.02, 0.995, 0.96), tmp_dir, "qplot")
        add_picture_fit(s, qplot, 0.52, 1.37, 6.40, 3.64)
        metric(s, "24.36%", "SR-TARGET (+ FIXED LR) REL. L2", 7.15, 1.55, 2.25,
               color=BLUE, fill="E7F2F8")
        metric(s, "91.93%", "DIRECT FULL-H REL. L2", 7.15, 2.53, 2.25,
               color=ORANGE, fill=PALE_ORANGE)
        box(s, 7.15, 3.56, 2.25, 1.16, fill=NAVY, radius=True)
        text(s, "Γ is nearly tied", 7.31, 3.73, 1.93, 0.23, size=12,
             color=MINT, bold=True, align=PP_ALIGN.CENTER)
        text(s, "The separation appears at nonzero q.", 7.31, 4.08, 1.93, 0.44,
             size=11, color=WHITE, align=PP_ALIGN.CENTER)
        footer(s, 5)
        add_notes(s, "Cartesian EPC exposed a finite-q failure in direct full-H", """
2:45–3:35. Across the full 2×2×2 Cartesian-AO tensor, the SR checkpoint is at
24.36 percent relative L2 while direct full-H is at 91.93 percent. At Gamma
they are essentially tied: 6.07 versus 5.82 percent. The failure appears at
nonzero q, where direct full-H both overestimates the coupling strength and
has much larger relative error. This is the central experimental result.
""")

        # 6 — pipeline validation including Au
        s = new_slide(prs)
        title(s, "Validation", "The machinery is already tested: gold first, then Cartesian MgO EPC", 6)
        box(s, 0.62, 1.47, 3.26, 3.49, fill=LIGHT, line=LINE, radius=True)
        status_chip(s, "Complete", 0.84, 1.70, 0.88, color=GREEN)
        text(s, "Bulk gold", 0.84, 2.14, 2.72, 0.36, size=21,
             color=NAVY, font="Cambria", bold=True)
        text(s, "• MACE-H machinery test completed\n\n• Nonpolar metallic case isolates implementation issues\n\n• Established a trusted starting point before MgO",
             0.84, 2.67, 2.75, 1.74, size=12, color=GRAY,
             line_spacing=1.05)
        conv = crop_asset(HERE / "epc_04_finite_difference_convergence.png",
                          (0.00, 0.02, 1.00, 0.97), tmp_dir, "fd")
        add_picture_fit(s, conv, 4.08, 1.47, 5.30, 2.85)
        metric(s, "0.00656%", "DFT STEP-SIZE CHANGE", 4.18, 4.48, 1.62,
               color=GREEN, fill="E4F4EE")
        metric(s, "2.9×10⁻¹¹", "WORST HERMITICITY / PEAK", 5.96, 4.48, 1.86,
               color=TEAL)
        metric(s, "301,056", "COMPLEX EPC COMPONENTS", 7.98, 4.48, 1.40,
               color=BLUE, fill="E7F2F8")
        footer(s, 6)
        add_notes(s, "The machinery is already tested: gold first, then Cartesian MgO EPC", """
3:35–4:20. I am not starting from an untested code path. Bulk gold exercised
the model machinery first. For MgO, the Cartesian EPC reference uses twelve
converged ABACUS calculations and contains 301,056 complex components. Doubling
the DFT displacement changes the tensor by only 0.00656 percent, and the
Hermiticity residual is 2.9 times 10 to the minus 11 of the tensor peak. These
checks support the finite-difference and indexing pipeline.
""")

        # 7 — architecture audit
        s = new_slide(prs)
        title(s, "Audit", "The previous “full-H” run was not the intended long-range model", 7)
        audit = crop_asset(HERE / "full_h_pipeline_audit_equilibrium.png",
                           (0.01, 0.03, 0.995, 0.97), tmp_dir, "audit")
        add_picture_fit(s, audit, 0.55, 1.38, 5.85, 3.33)
        rows = [
            ("Direct Full-H", "predicts total H", ORANGE),
            ("SR checkpoint", "predicts HSR", BLUE),
            ("Z* head", "absent", RED),
            ("ε∞ head", "absent", RED),
            ("LR wrapper", "fixed DFPT tensors", GREEN),
        ]
        y0 = 1.54
        for i, (a, b, c) in enumerate(rows):
            y = y0 + i * 0.55
            box(s, 6.66, y, 2.72, 0.43, fill=LIGHT, line=LINE, radius=True)
            text(s, a, 6.80, y + 0.09, 1.02, 0.20, size=9.8, color=c, bold=True)
            text(s, b, 7.87, y + 0.09, 1.34, 0.20, size=9.8, color=INK,
                 align=PP_ALIGN.RIGHT)
        box(s, 6.66, 4.42, 2.72, 0.64, fill=NAVY, radius=True)
        text(s, "The learned-tensor case never existed in these checkpoints.",
             6.82, 4.51, 2.40, 0.42, size=10, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        footer(s, 7)
        add_notes(s, "The previous “full-H” run was not the intended long-range model", """
4:20–5:05. The audit is decisive. Both checkpoints have the same parameter
count and neither contains a Born-charge or dielectric head. Direct full-H is
an independently trained total-H predictor. The SR model gets fixed DFPT
tensors through an external wrapper. Therefore the previous training did not
test the proposed learned-tensor architecture. That is the correction I am
making now.
""")

        # 8 — decomposition / interpretation
        s = new_slide(prs)
        title(s, "Controlled decomposition", "The SR advantage is not yet an analytic long-range result", 8)
        top_left = crop_asset(HERE / "epc_pipeline_decomposition.png",
                              (0.01, 0.02, 0.50, 0.50), tmp_dir, "decomp_main")
        bottom_right = crop_asset(HERE / "epc_pipeline_decomposition.png",
                                  (0.50, 0.51, 0.99, 0.99), tmp_dir, "decomp_lr")
        add_picture_fit(s, top_left, 0.55, 1.42, 4.32, 3.26)
        add_picture_fit(s, bottom_right, 5.13, 1.42, 4.32, 3.26)
        box(s, 0.72, 4.76, 8.56, 0.36, fill=NAVY, radius=True)
        rich_text(s, [
            ("A → B: ", MINT, True, 11),
            ("−0.000078 percentage points", WHITE, True, 11),
            ("   ·   analytic LR / needed residual: ", MID, False, 11),
            ("0.001606%", "F5B65B", True, 11),
        ], 0.92, 4.84, 8.16, 0.20, align=PP_ALIGN.CENTER)
        footer(s, 8)
        add_notes(s, "The SR advantage is not yet an analytic long-range result", """
5:05–5:45. The controlled A/B/D decomposition prevents an incorrect causal
claim. A is the SR checkpoint alone, B adds fixed analytic LR, and D is direct
full-H. A and B are visually and numerically identical on this 2×2×2 grid.
The analytic correction is only 0.001606 percent of the residual needed to
match DFT. The current advantage belongs to the SR checkpoint, not yet to the
analytic correction.
""")

        # 9 — response scan diagnosis
        s = new_slide(prs)
        title(s, "Diagnosis", "A model can match H while learning the wrong local slope", 9)
        scan_a = crop_asset(HERE / "response_scan_four_bucket_summary.png",
                            (0.02, 0.02, 0.50, 0.51), tmp_dir, "scan_a")
        scan_b = crop_asset(HERE / "response_scan_four_bucket_summary.png",
                            (0.50, 0.02, 0.99, 0.51), tmp_dir, "scan_b")
        add_picture_fit(s, scan_a, 0.55, 1.43, 4.15, 3.18)
        add_picture_fit(s, scan_b, 4.90, 1.43, 4.55, 3.18)
        rich_text(s, [
            ("Direct full-H raw H MAE: ", GRAY, False, 10.5),
            ("0.785 meV", ORANGE, True, 10.5),
            ("  vs SR ", GRAY, False, 10.5),
            ("0.842 meV", BLUE, True, 10.5),
            ("   ·   Full/SR slope RMSE: ", GRAY, False, 10.5),
            ("1.67×", RED, True, 11.5),
        ], 0.72, 4.72, 8.56, 0.25, align=PP_ALIGN.CENTER)
        text(s, "Matched-seed repeats or derivative-aware supervision are needed to separate target design from training variance.",
             0.72, 5.03, 8.56, 0.19, size=8.8, color=GRAY,
             align=PP_ALIGN.CENTER, italic=True)
        footer(s, 9)
        add_notes(s, "A model can match H while learning the wrong local slope", """
5:45–6:45. A continuous 25-point Mg-x scan explains the mismatch. Direct
full-H has slightly better raw Hamiltonian MAE, but its central-slope RMSE is
1.67 times worse. In 24.1 percent of matrix elements it is closer in H but
worse in slope. The mechanism is plausible, but because these are independent
training runs, matched-seed repeats or derivative-aware supervision are still
needed to distinguish target design from ordinary training variance.
""")

        # 10 — corrected architecture
        s = new_slide(prs)
        title(s, "Correction in progress", "Train HSR, Z*, and ε∞ together—then reconstruct full H", 10)
        # shared backbone
        box(s, 0.55, 2.12, 1.45, 1.10, fill=LIGHT, line=LINE, radius=True)
        text(s, "Structure", 0.70, 2.36, 1.15, 0.27, size=17,
             color=NAVY, font="Cambria", bold=True, align=PP_ALIGN.CENTER)
        text(s, "R, lattice", 0.70, 2.75, 1.15, 0.20, size=9.5,
             color=GRAY, align=PP_ALIGN.CENTER)
        line(s, 2.02, 2.67, 2.33, 2.67, color=TEAL, width=2.4, arrow=True)
        box(s, 2.39, 1.84, 1.42, 1.66, fill=NAVY, radius=True)
        text(s, "MACE-H", 2.52, 2.21, 1.16, 0.32, size=18, color=WHITE,
             font="Cambria", bold=True, align=PP_ALIGN.CENTER)
        text(s, "shared equivariant\nbackbone", 2.52, 2.70, 1.16, 0.42,
             size=9.5, color=MID, align=PP_ALIGN.CENTER)
        heads = [
            (4.18, 1.35, "HSR head", "local edge output", BLUE, "E7F2F8"),
            (4.18, 2.25, "Z* head", "per atom · ASR", ORANGE, PALE_ORANGE),
            (4.18, 3.15, "ε∞ head", "global · SPD", GREEN, "E4F4EE"),
        ]
        for x, y, head, body, color, fill in heads:
            line(s, 3.83, 2.67, x - 0.07, y + 0.35, color=TEAL, width=1.7,
                 arrow=True)
            box(s, x, y, 1.58, 0.72, fill=fill, line=LINE, radius=True)
            text(s, head, x + 0.08, y + 0.11, 1.42, 0.22, size=13.5,
                 color=color, font="Cambria", bold=True, align=PP_ALIGN.CENTER)
            text(s, body, x + 0.08, y + 0.42, 1.42, 0.16, size=8.5,
                 color=GRAY, align=PP_ALIGN.CENTER)
        line(s, 5.83, 2.60, 6.18, 2.60, color=TEAL, width=2.0, arrow=True)
        box(s, 6.24, 1.90, 1.48, 1.38, fill="E4F4EE", line=LINE, radius=True)
        text(s, "Analytic HLR", 6.36, 2.22, 1.24, 0.26, size=15,
             color=GREEN, font="Cambria", bold=True, align=PP_ALIGN.CENTER)
        text(s, "differentiable\nreconstruction", 6.36, 2.66, 1.24, 0.34,
             size=9, color=GRAY, align=PP_ALIGN.CENTER)
        line(s, 7.76, 2.60, 8.05, 2.60, color=TEAL, width=2.3, arrow=True)
        box(s, 8.11, 1.90, 1.31, 1.38, fill=NAVY, radius=True)
        text(s, "Hfull", 8.23, 2.25, 1.07, 0.28, size=19, color=WHITE,
             font="Cambria", bold=True, align=PP_ALIGN.CENTER)
        text(s, "+ Cartesian EPC", 8.23, 2.72, 1.07, 0.20, size=8.8,
             color=MINT, bold=True, align=PP_ALIGN.CENTER)
        box(s, 0.65, 4.32, 8.70, 0.66, fill=LIGHT, line=LINE, radius=True)
        rich_text(s, [
            ("Joint objective:  ", NAVY, True, 11),
            ("LHSR", BLUE, True, 11),
            (" + ", GRAY, False, 11),
            ("LZ*", ORANGE, True, 11),
            (" + ", GRAY, False, 11),
            ("Lε∞", GREEN, True, 11),
            (" + ", GRAY, False, 11),
            ("LHfull reconstruction", NAVY, True, 11),
            (" + response consistency", RED, True, 11),
        ], 0.88, 4.52, 8.24, 0.23, align=PP_ALIGN.CENTER)
        footer(s, 10)
        add_notes(s, "Train HSR, Z*, and ε∞ together—then reconstruct full H", """
6:45–8:00. This is the corrected implementation. A shared MACE-H backbone
feeds three outputs: the short-range Hamiltonian, an atom-resolved Born-charge
head with the acoustic sum rule enforced, and a global symmetric
positive-definite electronic dielectric head. Those tensors feed a
differentiable analytic long-range reconstruction. The training objective
includes direct tensor supervision and full-H reconstruction; if schedule
allows, I also want a finite-difference response-consistency term. The key
ablation is A/B/C/D: SR only, fixed tensors, learned tensors, and direct full-H.
""")

        # 11 — timeline and handoff
        s = new_slide(prs)
        title(s, "Next steps + handoff", "What I will finish before August 20—and what the group should decide", 11)
        timeline = [
            ("13", "ALIGN", "objective + success gates", TEAL),
            ("14–15", "LABELS", "54-atom Z*/ε∞ fast-vs-anchor gate", ORANGE),
            ("15–17", "BUILD", "heads, loader, constraints, tests", BLUE),
            ("17–18", "TRAIN", "joint smoke + matched-seed ablations", RED),
            ("18–19", "EVALUATE", "tensor → H → Cartesian EPC", GREEN),
            ("20", "HANDOFF", "refactor, provenance, runbook", NAVY),
        ]
        x0, step = 0.62, 1.47
        line(s, 0.86, 2.12, 8.98, 2.12, color=LINE, width=3.0)
        for i, (day, stage, detail, color) in enumerate(timeline):
            x = x0 + i * step
            box(s, x + 0.14, 1.93, 0.38, 0.38, fill=color, radius=True)
            text(s, day, x - 0.02, 1.47, 0.70, 0.25, size=10.5,
                 color=color, bold=True, align=PP_ALIGN.CENTER)
            text(s, stage, x - 0.17, 2.48, 1.08, 0.20, size=8.7,
                 color=color, bold=True, align=PP_ALIGN.CENTER)
            text(s, detail, x - 0.28, 2.82, 1.30, 0.62, size=8.8,
                 color=GRAY, align=PP_ALIGN.CENTER)
        box(s, 0.62, 3.76, 4.18, 1.17, fill=LIGHT, line=LINE, radius=True)
        text(s, "LEAVE BEHIND", 0.84, 3.96, 1.20, 0.18, size=8.5,
             color=TEAL, bold=True)
        text(s, "Clean physics interfaces · tensor-head configs · A/B/C/D evaluation · tests · frozen provenance · restartable runbook",
             0.84, 4.24, 3.72, 0.49, size=10.5, color=NAVY, bold=True)
        box(s, 5.02, 3.76, 4.36, 1.17, fill=NAVY, radius=True)
        text(s, "GROUP DECISIONS BEFORE AUG 20", 5.24, 3.96, 2.48, 0.18,
             size=8.5, color=MINT, bold=True)
        text(s, "Paper result or reusable platform?  ·  MgO-only or multi-material?  ·  Denser-q EPC or more tensor labels?  ·  Who owns the next campaign?",
             5.24, 4.24, 3.90, 0.49, size=10.5, color=WHITE, bold=True)
        footer(s, 11)
        add_notes(s, "What I will finish before August 20—and what the group should decide", """
8:00–9:15. The critical path is: validate the fast tensor-label profile against
the anchor, implement and test the two equivariant heads, run a joint-training
smoke and matched-seed ablations, then evaluate in the order tensor accuracy,
Hamiltonian reconstruction, Cartesian EPC. By August 20 I will leave clean
interfaces, configs, tests, provenance, and a restartable runbook. I need the
group to decide whether the continuation is primarily a paper result or a
reusable platform, whether scope stays MgO-only, and whether the next compute
budget goes to denser q or more tensor-labelled structures.
""")

        # 12 — takeaways
        s = new_slide(prs, NAVY)
        box(s, 0.0, 0.0, 0.18, H, fill=MINT)
        text(s, "Takeaways", 0.75, 0.62, 8.5, 0.55, size=32, color=WHITE,
             font="Cambria", bold=True)
        takeaways = [
            ("1", "The model machinery has been tested on bulk gold.", GREEN),
            ("2", "Cartesian MgO EPC works—and revealed a finite-q response failure.", BLUE),
            ("3", "The old direct full-H run did not contain Z* or ε∞ heads.", ORANGE),
            ("4", "The correction is a composed model: HSR + learned tensors + analytic HLR.", MINT),
            ("5", "By Aug 20: tested MVP, refactored code, reproducible handoff, and an agreed next owner/scope.", WHITE),
        ]
        for i, (num, body, color) in enumerate(takeaways):
            y = 1.52 + i * 0.68
            box(s, 0.76, y, 0.36, 0.36, fill=color if color != WHITE else MID,
                radius=True)
            text(s, num, 0.76, y + 0.05, 0.36, 0.18, size=10.5,
                 color=NAVY, bold=True, align=PP_ALIGN.CENTER)
            text(s, body, 1.35, y - 0.01, 7.80, 0.42, size=13.5,
                 color=PALE_BLUE if color != WHITE else WHITE,
                 bold=(i == 4), valign=MSO_ANCHOR.MIDDLE)
        text(s, "Questions  ·  priorities for the final week?", 0.75, 5.05, 8.5, 0.27,
             size=13, color=MINT, bold=True)
        add_notes(s, "Takeaways", """
9:15–10:00. The machinery and Cartesian EPC pipeline are complete. The key
learning is that the previous direct full-H training was not the proposed
long-range architecture. I am correcting that now with explicit Born-charge
and dielectric heads plus analytic reconstruction. Before I leave, I want to
deliver a tested minimum viable implementation and a handoff aligned with the
group's preferred scientific direction. Then stop and invite questions.
""")

        # 13 — backup scope
        s = new_slide(prs)
        title(s, "Backup · scope", "What the current EPC tensor does—and does not—establish", 13)
        direction = crop_asset(HERE / "epc_results_06_angular_direction_accuracy.png",
                               (0.01, 0.02, 0.995, 0.97), tmp_dir, "direction")
        add_picture_fit(s, direction, 0.55, 1.42, 5.85, 3.42)
        box(s, 6.68, 1.43, 2.70, 3.41, fill=LIGHT, line=LINE, radius=True)
        status_chip(s, "Established", 6.89, 1.66, 1.08, color=GREEN)
        text(s, "• Cartesian AO dH/dτ\n• 2×2×2 k/q grid\n• 2 atoms × 3 directions\n• finite-difference + Hermiticity checks",
             6.89, 2.08, 2.24, 1.18, size=11, color=INK)
        status_chip(s, "Not yet", 6.89, 3.43, 0.82, color=ORANGE)
        text(s, "• band/mode-resolved EPC\n• phonons + mass/frequency factors\n• dS/dτ contribution\n• converged small-q behavior",
             6.89, 3.83, 2.24, 0.84, size=10.5, color=GRAY)
        footer(s, 13)
        add_notes(s, "What the current EPC tensor does—and does not—establish", """
Backup. The current result is a Cartesian atomic-orbital Hamiltonian
derivative. It does not yet include phonon eigenvectors and frequencies,
electronic eigenvectors, mass factors, or the overlap derivative, and the
2×2×2 q grid is not a converged small-q study. The direction metric still
shows a consistent checkpoint-level difference across Cartesian components.
""")

        # 14 — backup A/B/C/D
        s = new_slide(prs)
        title(s, "Backup · experiment design", "The comparison that will isolate the learned-tensor contribution", 14)
        cols = [
            ("A", "SR only", "Current", "dHSR/dτ", BLUE),
            ("B", "SR + fixed LR", "Current", "uses DFPT Z*, ε∞", GREEN),
            ("C", "SR + learned LR", "Next", "predicted Z*, ε∞", ORANGE),
            ("D", "Direct full-H", "Baseline", "independent total-H fit", RED),
        ]
        xs = [0.58, 2.88, 5.18, 7.48]
        for (letter, name, state, desc, color), x in zip(cols, xs):
            box(s, x, 1.55, 1.94, 2.50, fill=LIGHT, line=LINE, radius=True)
            box(s, x + 0.17, 1.77, 0.38, 0.38, fill=color, radius=True)
            text(s, letter, x + 0.17, 1.82, 0.38, 0.18, size=11,
                 color=WHITE, bold=True, align=PP_ALIGN.CENTER)
            text(s, name, x + 0.17, 2.36, 1.60, 0.46, size=16,
                 color=NAVY, font="Cambria", bold=True, align=PP_ALIGN.CENTER)
            status_chip(s, state, x + 0.51, 3.01, 0.92, color=color)
            text(s, desc, x + 0.17, 3.48, 1.60, 0.30, size=9.5,
                 color=GRAY, align=PP_ALIGN.CENTER)
        box(s, 0.75, 4.41, 8.50, 0.60, fill=NAVY, radius=True)
        text(s, "B − A isolates fixed analytic LR   ·   C − A isolates learned tensors   ·   C vs D tests the composed objective",
             0.95, 4.61, 8.10, 0.20, size=11, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER)
        footer(s, 14)
        add_notes(s, "The comparison that will isolate the learned-tensor contribution", """
Backup. A and B exist now and show that fixed LR is negligible on this grid.
D is the current direct full-H baseline. C is the missing learned-tensor case.
B minus A isolates the fixed analytic contribution; C minus A measures the
effect of learned structure-dependent tensors; C versus D answers whether the
composed objective improves the physical response.
""")

        # 15 — backup risks and decisions
        s = new_slide(prs)
        title(s, "Backup · execution", "If time compresses, preserve the causal test and the handoff", 15)
        headers = ["Priority", "Must finish", "Can defer"]
        widths = [1.55, 3.60, 3.78]
        xs = [0.58, 2.13, 5.73]
        for h, x, w in zip(headers, xs, widths):
            box(s, x, 1.48, w, 0.48, fill=NAVY if h != "Can defer" else TEAL,
                radius=True)
            text(s, h, x + 0.12, 1.61, w - 0.24, 0.19, size=10.5,
                 color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        rows = [
            ("1", "Tensor-label validity + constrained heads", "Large tensor-label campaign"),
            ("2", "A/B/C/D matched comparison", "Band/mode projection"),
            ("3", "Restartable code, tests, provenance", "Converged dense-q production"),
        ]
        for i, (p, must, defer) in enumerate(rows):
            y = 2.10 + i * 0.88
            box(s, 0.58, y, 1.55, 0.72, fill=LIGHT, line=LINE)
            box(s, 2.13, y, 3.60, 0.72, fill=WHITE, line=LINE)
            box(s, 5.73, y, 3.78, 0.72, fill=WHITE, line=LINE)
            text(s, p, 0.58, y + 0.21, 1.55, 0.24, size=14, color=TEAL,
                 bold=True, align=PP_ALIGN.CENTER)
            text(s, must, 2.31, y + 0.17, 3.24, 0.34, size=11.5,
                 color=NAVY, bold=True, align=PP_ALIGN.CENTER,
                 valign=MSO_ANCHOR.MIDDLE)
            text(s, defer, 5.93, y + 0.17, 3.38, 0.34, size=11,
                 color=GRAY, align=PP_ALIGN.CENTER,
                 valign=MSO_ANCHOR.MIDDLE)
        box(s, 0.72, 4.98, 8.64, 0.25, fill=LIGHT, radius=True)
        text(s, "Minimum viable handoff: one correct learned-tensor checkpoint + an experiment someone else can reproduce.",
             0.85, 5.03, 8.38, 0.15, size=9.3, color=NAVY, bold=True,
             align=PP_ALIGN.CENTER)
        footer(s, 15)
        add_notes(s, "If time compresses, preserve the causal test and the handoff", """
Backup. If the week compresses, the minimum viable outcome is not the largest
campaign. It is valid tensor labels, constrained heads, one matched A/B/C/D
comparison, and a reproducible handoff. Band- and mode-resolved EPC, a large
tensor dataset, and dense-q production can follow after ownership is clear.
""")

        prs.save(str(OUTPUT))

    note_lines = [
        "# MACE-H-LR project update — speaker notes",
        "",
        "Presentation date: August 13, 2026  ",
        "Target length: 10 minutes  ",
        "Slides 13–15 are backup.",
        "",
    ]
    for i, (slide_title, slide_notes) in enumerate(SLIDE_NOTES, 1):
        note_lines.extend([f"## {i}. {slide_title}", "", slide_notes, ""])
    NOTES_OUTPUT.write_text("\n".join(note_lines), encoding="utf-8")
    print(OUTPUT)
    print(NOTES_OUTPUT)


if __name__ == "__main__":
    build()
